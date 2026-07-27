import fitz
from pptx import Presentation
from pptx.util import Pt
import os

def pdf_to_editable_ppt(pdf_path, pptx_path):
    doc = fitz.open(pdf_path)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    
    # Standard pptx width is 10 inches, height is 7.5 inches.
    # Typical PDF is 8.5 x 11 inches. We might need a scale factor.
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # Scale factor mapping PDF dimensions to PPTX dimensions
        # PPTX default slide: 9144000x6858000 EMU (10x7.5 inches) -> 720x540 points
        pdf_width = page.rect.width
        pdf_height = page.rect.height
        
        scale_x = prs.slide_width / Pt(pdf_width) if pdf_width else 1
        scale_y = prs.slide_height / Pt(pdf_height) if pdf_height else 1

        slide = prs.slides.add_slide(blank_slide_layout)
        
        blocks = page.get_text("blocks")
        for b in blocks:
            x0, y0, x1, y1, text, block_no, block_type = b
            if block_type == 0: # text
                # clean up text
                text = text.strip()
                if not text:
                    continue
                
                # Apply scale
                left = Pt(x0) * scale_x
                top = Pt(y0) * scale_y
                width = Pt(x1 - x0) * scale_x
                height = Pt(y1 - y0) * scale_y
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = text
                # We can set word wrapping to true
                tf.word_wrap = True
                
            elif block_type == 1: # image
                # getting image is a bit harder from block directly if it's complex,
                # but we can try extracting it from page
                try:
                    # just skip for now in test, or extract via doc.extract_image
                    pass
                except:
                    pass

    doc.close()
    prs.save(pptx_path)

if __name__ == "__main__":
    pdf_to_editable_ppt("test_text.pdf", "test_ppt.pptx")
    print("Done")
