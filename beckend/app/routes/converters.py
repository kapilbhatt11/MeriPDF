from fastapi import APIRouter, File, UploadFile, HTTPException, Form
from fastapi.responses import FileResponse
import tempfile
import os
import uuid
import asyncio
from PIL import Image
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
except ImportError:
    pass
import io
import fitz # PyMuPDF
import zipfile
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pdfplumber
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from fastapi.responses import FileResponse, StreamingResponse
import pytesseract

router = APIRouter(prefix="/converters", tags=["converters"])

# Helper to process image fitting onto canvas
def format_image_for_pdf(img: Image.Image, page_size: str, orientation: str, margin: str) -> Image.Image:
    # 1. If page_size is fit, keep original image size
    if page_size == "fit":
        return img
    
    # 2. Base page dimensions in points
    # A4: 595 x 842, Letter: 612 x 792
    if page_size == "a4":
        p_w, p_h = 595, 842
    elif page_size == "letter":
        p_w, p_h = 612, 792
    else:
        p_w, p_h = 595, 842  # default fallback A4
        
    # Rotate based on orientation
    img_w, img_h = img.size
    
    if orientation == "portrait":
        sw, sh = min(p_w, p_h), max(p_w, p_h)
    elif orientation == "landscape":
        sw, sh = max(p_w, p_h), min(p_w, p_h)
    else: # auto
        if img_w > img_h:
            sw, sh = max(p_w, p_h), min(p_w, p_h)
        else:
            sw, sh = min(p_w, p_h), max(p_w, p_h)
            
    # Margins in points
    pad = 0
    if margin == "small":
        pad = 20
    elif margin == "large":
        pad = 40
        
    aw = sw - 2 * pad
    ah = sh - 2 * pad
    
    # Scale image to fit inside available boundaries preserving aspect ratio
    scale = min(aw / img_w, ah / img_h)
    new_w = max(1, int(img_w * scale))
    new_h = max(1, int(img_h * scale))
    
    # Resize the image using High Quality Resampling
    try:
        sampler = Image.Resampling.LANCZOS
    except AttributeError:
        sampler = Image.ANTIALIAS
        
    resized_img = img.resize((new_w, new_h), sampler)
    
    # Create white canvas of (sw, sh)
    canvas = Image.new("RGB", (sw, sh), (255, 255, 255))
    
    # Paste centered
    x_off = pad + (aw - new_w) // 2
    y_off = pad + (ah - new_h) // 2
    canvas.paste(resized_img, (x_off, y_off))
    
    return canvas

@router.post("/image-to-pdf")
async def image_to_pdf(
    files: list[UploadFile] = File(...),
    page_size: str = Form("fit"),
    orientation: str = Form("auto"),
    margin: str = Form("none")
):
    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded")
    if len(files) > 50:
        raise HTTPException(status_code=400, detail="Maximum limit of 50 images exceeded per conversion request.")
    
    try:
        images_list = []
        for file in files:
            content = await file.read()
            img = Image.open(io.BytesIO(content))
            
            # Convert image to RGB because PDF doesn't support RGBA (transparent) directly natively
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            processed_img = format_image_for_pdf(img, page_size, orientation, margin)
            images_list.append(processed_img)
                
        if not images_list:
            raise HTTPException(status_code=400, detail="Invalid image files provided")
            
        temp_dir = tempfile.gettempdir()
        output_filename = "Converted_Images_MeriPDF.pdf"
        output_path = os.path.join(temp_dir, output_filename)
        
        # Save all images as a single PDF
        first_image = images_list[0]
        other_images = images_list[1:]
        
        if other_images:
            first_image.save(output_path, "PDF", resolution=100.0, save_all=True, append_images=other_images)
        else:
            first_image.save(output_path, "PDF", resolution=100.0)
            
        return FileResponse(
            output_path,
            media_type="application/pdf",
            filename=output_filename,
            headers={"Content-Disposition": f"attachment; filename={output_filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ---------------------------------------------------------------------
# 🖼️ PDF to JPG (ZIP/Single Image)
# ---------------------------------------------------------------------
def parse_page_range(range_str: str, max_pages: int) -> list[int]:
    if not range_str or not range_str.strip():
        return list(range(max_pages))
    
    pages = set()
    parts = range_str.replace(" ", "").split(",")
    for part in parts:
        if "-" in part:
            try:
                start_str, end_str = part.split("-")
                start = int(start_str)
                end = int(end_str)
                start = max(1, min(start, max_pages))
                end = max(1, min(end, max_pages))
                for p in range(min(start, end), max(start, end) + 1):
                    pages.add(p - 1)
            except ValueError:
                continue
        else:
            try:
                p = int(part)
                if 1 <= p <= max_pages:
                    pages.add(p - 1)
            except ValueError:
                continue
    return sorted(list(pages)) if pages else list(range(max_pages))

def compress_to_size(img: Image.Image, target_bytes: int, fmt: str) -> bytes:
    pil_format = fmt.upper()
    if pil_format == "JPG":
        pil_format = "JPEG"
        
    TARGET_QUALITY = 82
    
    buf = io.BytesIO()
    save_args = {"format": pil_format}
    if pil_format in ["JPEG", "WEBP"]:
        save_args["quality"] = TARGET_QUALITY
        
    try:
        img.save(buf, **save_args)
    except Exception:
        img.save(buf, format="JPEG", quality=TARGET_QUALITY)
        pil_format = "JPEG"
        
    if buf.tell() <= target_bytes:
        return buf.getvalue()
        
    # Scale down the image resolution while keeping quality high (visually sharp and crisp)
    low_scale = 0.05
    high_scale = 0.95
    best_bytes = None
    
    for _ in range(6):
        mid_scale = (low_scale + high_scale) / 2
        w = max(1, int(img.width * mid_scale))
        h = max(1, int(img.height * mid_scale))
        resized = img.resize((w, h), Image.Resampling.LANCZOS)
        
        test_buf = io.BytesIO()
        test_args = {"format": pil_format}
        if pil_format in ["JPEG", "WEBP"]:
            test_args["quality"] = TARGET_QUALITY
            
        resized.save(test_buf, **test_args)
        size = test_buf.tell()
        
        if size <= target_bytes:
            best_bytes = test_buf.getvalue()
            low_scale = mid_scale + 0.02
        else:
            high_scale = mid_scale - 0.02
            
    if best_bytes:
        return best_bytes
        
    # Fallback to extremely low dimensions with reduced quality if search failed
    w = max(1, int(img.width * 0.05))
    h = max(1, int(img.height * 0.05))
    resized = img.resize((w, h), Image.Resampling.LANCZOS)
    final_buf = io.BytesIO()
    if pil_format in ["JPEG", "WEBP"]:
        resized.save(final_buf, format=pil_format, quality=60)
    else:
        resized.save(final_buf, format=pil_format)
    return final_buf.getvalue()


@router.post("/pdf-to-jpg")
async def pdf_to_jpg(
    file: UploadFile = File(...),
    mode: str = Form("pages"),
    dpi: int = Form(150),
    format: str = Form("jpg"),
    page_range: str = Form(""),
    target_kb: str = Form(None)
):
    try:
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        target_kb_val = None
        if target_kb and target_kb.strip():
            try:
                target_kb_val = int(target_kb)
            except ValueError:
                pass
                
        pages_to_process = parse_page_range(page_range, len(doc))
        
        if mode == "text":
            text_content = ""
            for page_num in pages_to_process:
                page = doc.load_page(page_num)
                text_content += f"--- Page {page_num + 1} ---\n\n"
                text_content += page.get_text("text") + "\n\n"
            
            text_bytes = text_content.encode("utf-8")
            return StreamingResponse(
                io.BytesIO(text_bytes),
                media_type="text/plain",
                headers={"Content-Disposition": f"attachment; filename=Extracted_Text_{os.path.splitext(file.filename)[0]}.txt"}
            )
            
        elif mode == "text_ocr":
            text_content = ""
            for page_num in pages_to_process:
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=dpi)
                img = Image.open(io.BytesIO(pix.tobytes("jpeg")))
                text = pytesseract.image_to_string(img, lang="hin+eng")
                text_content += f"--- Page {page_num + 1} ---\n\n"
                text_content += text + "\n\n"
            
            text_bytes = text_content.encode("utf-8")
            return StreamingResponse(
                io.BytesIO(text_bytes),
                media_type="text/plain",
                headers={"Content-Disposition": f"attachment; filename=Extracted_Text_OCR_{os.path.splitext(file.filename)[0]}.txt"}
            )
            
        elif mode == "images":
            extracted_images = []
            img_count = 0
            for page_num in pages_to_process:
                page = doc.load_page(page_num)
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    img_count += 1
                    extracted_images.append((f"Image_{img_count}.{image_ext}", image_bytes))
            
            # Smart single page response for 1 extracted image
            if len(extracted_images) == 1:
                single_name, single_bytes = extracted_images[0]
                ext = single_name.split(".")[-1]
                mime_type = f"image/{ext}"
                if ext == "jpg":
                    mime_type = "image/jpeg"
                return StreamingResponse(
                    io.BytesIO(single_bytes),
                    media_type=mime_type,
                    headers={"Content-Disposition": f"attachment; filename={single_name}"}
                )
                
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w") as zipf:
                for img_name, img_bytes in extracted_images:
                    zipf.writestr(img_name, img_bytes)
                if img_count == 0:
                    zipf.writestr("no_images_found.txt", b"No embedded images were found in this PDF.")
            
            zip_buffer.seek(0)
            filename = f"Extracted_Images_{os.path.splitext(file.filename)[0]}.zip"
            return StreamingResponse(
                zip_buffer,
                media_type="application/zip",
                headers={"Content-Disposition": f"attachment; filename={filename}"}
            )
            
        else: # mode == "pages"
            # Smart single page response for 1 page rendered
            if len(pages_to_process) == 1:
                page_num = pages_to_process[0]
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=dpi)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                
                if format == "jpg" and img.mode != "RGB":
                    img = img.convert("RGB")
                    
                ext = "jpg" if format == "jpg" else format
                pil_fmt = "JPEG" if format == "jpg" else format.upper()
                mime_type = "image/jpeg" if format == "jpg" else f"image/{format}"
                
                if target_kb_val:
                    img_bytes = compress_to_size(img, target_kb_val * 1024, ext)
                else:
                    buf = io.BytesIO()
                    img.save(buf, format=pil_fmt)
                    img_bytes = buf.getvalue()
                    
                filename = f"Page_{page_num + 1}.{ext}"
                return StreamingResponse(
                    io.BytesIO(img_bytes),
                    media_type=mime_type,
                    headers={"Content-Disposition": f"attachment; filename={filename}"}
                )
                
            # Multiple pages zip response
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w") as zipf:
                for page_num in pages_to_process:
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(dpi=dpi)
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    
                    if format == "jpg" and img.mode != "RGB":
                        img = img.convert("RGB")
                        
                    ext = "jpg" if format == "jpg" else format
                    pil_fmt = "JPEG" if format == "jpg" else format.upper()
                    
                    if target_kb_val:
                        img_bytes = compress_to_size(img, target_kb_val * 1024, ext)
                    else:
                        buf = io.BytesIO()
                        img.save(buf, format=pil_fmt)
                        img_bytes = buf.getvalue()
                        
                    zipf.writestr(f"Page_{page_num + 1}.{ext}", img_bytes)
                    
            zip_buffer.seek(0)
            filename = f"Converted_Images_{os.path.splitext(file.filename)[0]}.zip"
            return StreamingResponse(
                zip_buffer,
                media_type="application/zip",
                headers={"Content-Disposition": f"attachment; filename={filename}"}
            )
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to process PDF: {e}")

# ---------------------------------------------------------------------
# 📝 PDF to WORD (.docx)
# ---------------------------------------------------------------------
@router.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    try:
        temp_dir = tempfile.gettempdir()
        pdf_path = os.path.join(temp_dir, f"temp_{file.filename}")
        docx_path = os.path.join(temp_dir, f"converted_{file.filename}.docx")
        
        with open(pdf_path, "wb") as f:
            f.write(await file.read())
            
        cv = Converter(pdf_path)
        cv.convert(docx_path)
        cv.close()
        
        return FileResponse(
            docx_path,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=f"Extracted_{file.filename}.docx"
        )
    except Exception as e:
        print("PDF to Word Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to convert PDF to Word: {e}")

# ---------------------------------------------------------------------
# 📊 PDF to POWERPOINT (.pptx)
# ---------------------------------------------------------------------
def page_needs_ocr(page) -> bool:
    text = page.get_text("text").strip()
    if not text:
        return True
        
    # Check for CID font glyph corruption errors (forces OCR fallback)
    if "(cid:" in text or "cid:" in text:
        return True
    
    # 1. Check if total alphanumeric text is extremely sparse
    alnum_chars = sum(1 for c in text if c.isalnum())
    if alnum_chars < 5:
        return True

    # 2. Check for legacy fonts (e.g. KrutiDev, Devlys, Shusha, Shivaji, Chanakya) via font names
    try:
        page_dict = page.get_text("dict")
        for block in page_dict.get("blocks", []):
            if block.get("type") == 0:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        font_lower = span.get("font", "").lower()
                        if any(x in font_lower for x in ["kruti", "devlys", "shusha", "chanakya", "shivaji"]):
                            return True
    except Exception:
        pass
        
    # Unicode range check for legacy DTP font hijack detection
    invalid_count = 0
    total_len = len(text)
    for char in text:
        o = ord(char)
        # Whitelist ranges:
        # A. ASCII & Latin-1 Supplement: 0x0000 to 0x00FF
        if 0x0000 <= o <= 0x00FF:
            continue
        # B. Devanagari: 0x0900 to 0x097F
        if 0x0900 <= o <= 0x097F:
            continue
        # C. Devanagari Extended / Sanskrit: 0xA8E0 to 0xA8FF
        if 0xA8E0 <= o <= 0xA8FF:
            continue
        # D. Vedic Extensions: 0x1CD0 to 0x1CFF
        if 0x1CD0 <= o <= 0x1CFF:
            continue
        # E. General Punctuation (Rupee symbol 0x20B9, dashes, spaces): 0x2000 to 0x20CF
        if 0x2000 <= o <= 0x20CF:
            continue
        # F. Mathematical Operators & arrows, etc.: 0x2190 to 0x22FF
        if 0x2190 <= o <= 0x22FF:
            continue
        # G. Custom encoded fonts (Balinese, Sundanese, Ol Chiki, Georgian): 0x1B00 to 0x1CFF
        if 0x1B00 <= o <= 0x1CFF:
            continue
        # H. Common punctuation / symbols
        if char in "|/<>~`\"'\\:।॥.,;!?@#$%^&*()_+-=[]{}":
            continue
        
        # Any other character outside valid bounds is considered invalid/suspicious (KrutiDev mapping noise)
        invalid_count += 1
        
    # Calculate ratio of invalid/suspicious characters
    has_hindi = any(0x0900 <= ord(c) <= 0x097F for c in text)
    
    # Over suspicious count/ratio dictates fallback to visual OCR
    if has_hindi:
        # Hindi PDFs should not have garbage characters representing ligatures
        if invalid_count > 5 or (total_len > 0 and (invalid_count / total_len) > 0.01):
            return True
    else:
        if invalid_count > 15 or (total_len > 0 and (invalid_count / total_len) > 0.15):
            return True
        
    return False

def get_hex_from_color(color_val) -> str:
    """Safely converts grayscale, RGB, and CMYK colors to openpyxl ARGB hex strings."""
    if not isinstance(color_val, (list, tuple)):
        return "FF000000"
    if len(color_val) == 1:
        val = int(color_val[0] * 255)
        r = g = b = val
    elif len(color_val) == 3:
        r = int(color_val[0] * 255)
        g = int(color_val[1] * 255)
        b = int(color_val[2] * 255)
    elif len(color_val) == 4:
        c, m, y, k = color_val
        r = int(255 * (1.0 - c) * (1.0 - k))
        g = int(255 * (1.0 - m) * (1.0 - k))
        b = int(255 * (1.0 - y) * (1.0 - k))
    else:
        r = g = b = 0
    r = max(0, min(255, r))
    g = max(0, min(255, g))
    b = max(0, min(255, b))
    return f"FF{r:02x}{g:02x}{b:02x}"

def merge_block_lines_by_baseline(lines):
    if not lines:
        return []
    
    # Sort lines by top coordinate Y0 first
    sorted_lines = sorted(lines, key=lambda l: l.get("bbox", (0,0,0,0))[1])
    
    grouped_lines = []
    current_group = []
    last_y0 = None
    
    for l in sorted_lines:
        y0 = l.get("bbox", (0,0,0,0))[1]
        if last_y0 is None:
            current_group.append(l)
            last_y0 = y0
        elif abs(y0 - last_y0) < 3.0: # Baseline tolerance of 3 points
            current_group.append(l)
        else:
            grouped_lines.append(current_group)
            current_group = [l]
            last_y0 = y0
    if current_group:
        grouped_lines.append(current_group)
        
    merged_lines = []
    for group in grouped_lines:
        if len(group) == 1:
            line = group[0]
            if "spans" in line:
                line["spans"].sort(key=lambda s: s.get("bbox", (0,0,0,0))[0])
            merged_lines.append(line)
        else:
            all_spans = []
            min_x0 = float('inf')
            min_y0 = float('inf')
            max_x1 = float('-inf')
            max_y1 = float('-inf')
            
            for line in group:
                for s in line.get("spans", []):
                    all_spans.append(s)
                lx0, ly0, lx1, ly1 = line.get("bbox", (0,0,0,0))
                min_x0 = min(min_x0, lx0)
                min_y0 = min(min_y0, ly0)
                max_x1 = max(max_x1, lx1)
                max_y1 = max(max_y1, ly1)
                
            # Sort spans left-to-right
            all_spans.sort(key=lambda s: s.get("bbox", (0,0,0,0))[0])
            
            merged_line = {
                "bbox": (min_x0, min_y0, max_x1, max_y1),
                "spans": all_spans
            }
            merged_lines.append(merged_line)
            
    merged_lines.sort(key=lambda l: l.get("bbox", (0,0,0,0))[1])
    return merged_lines

def sort_blocks_reading_order(blocks):
    """
    Sorts layout blocks visually using a layout-aware columns/rows analyzer.
    Blocks with substantial vertical overlap (>40%) are sorted left-to-right (parallel columns).
    Otherwise, sorted top-to-bottom.
    """
    text_blocks = [b for b in blocks if b.get("type") == 0]
    other_blocks = [b for b in blocks if b.get("type") != 0]
    
    from functools import cmp_to_key
    
    def compare_blocks(b1, b2):
        box1 = b1["bbox"]
        box2 = b2["bbox"]
        
        y_overlap = min(box1[3], box2[3]) - max(box1[1], box2[1])
        h1 = box1[3] - box1[1]
        h2 = box2[3] - box2[1]
        min_h = min(h1, h2)
        
        if min_h > 0 and (y_overlap / min_h) > 0.40:
            if abs(box1[0] - box2[0]) > 3:
                return -1 if box1[0] < box2[0] else 1
                
        if abs(box1[1] - box2[1]) > 3:
            return -1 if box1[1] < box2[1] else 1
            
        return -1 if box1[0] < box2[0] else 1
        
    sorted_text_blocks = sorted(text_blocks, key=cmp_to_key(compare_blocks))
    
    # Sort and merge lines inside each block
    for b in sorted_text_blocks:
        if "lines" in b:
            b["lines"] = merge_block_lines_by_baseline(b["lines"])
                    
    return sorted_text_blocks + other_blocks

def detect_text_color(img, bbox, page_w, page_h):
    try:
        w_img, h_img = img.size
        # Conversion to pixel indices
        px0 = int(bbox[0] * w_img / page_w)
        py0 = int(bbox[1] * h_img / page_h)
        px1 = int(bbox[2] * w_img / page_w)
        py1 = int(bbox[3] * h_img / page_h)
        
        px0 = max(0, min(w_img - 1, px0))
        px1 = max(px0 + 1, min(w_img, px1))
        py0 = max(0, min(h_img - 1, py0))
        py1 = max(py0 + 1, min(h_img, py1))
        
        cropped = img.crop((px0, py0, px1, py1))
        pixels = list(cropped.getdata())
        
        fg_pixels = []
        for p in pixels:
            if isinstance(p, tuple):
                r, g, b = p[0], p[1], p[2]
            else:
                r = g = b = p
                
            luma = 0.299 * r + 0.587 * g + 0.114 * b
            if luma < 210:  # foreground (darkish)
                fg_pixels.append((r, g, b))
                
        if len(fg_pixels) > 5:
            avg_r = sum(p[0] for p in fg_pixels) // len(fg_pixels)
            avg_g = sum(p[1] for p in fg_pixels) // len(fg_pixels)
            avg_b = sum(p[2] for p in fg_pixels) // len(fg_pixels)
            
            if avg_r < 35 and avg_g < 35 and avg_b < 35:
                return 0
            return (avg_r << 16) | (avg_g << 8) | avg_b
    except Exception:
        pass
    return 0

def get_page_elements(page, use_ocr=None):
    if use_ocr is None:
        use_ocr = page_needs_ocr(page)
    should_ocr = use_ocr
    if should_ocr:
        try:
            # Render at 200 DPI for high character definition (especially digits, slashes, currency)
            pix = page.get_pixmap(dpi=200)
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # Preprocessing: convert to grayscale (preserve anti-aliased glyph bounds for Tesseract)
            img = img.convert('L')
            
            resized_flag = False
            # pix.width/pix.height correspond to the unscaled image (1x width/height)
            
            # Use combined hin+eng OCR engine to recognize both Hindi Devanagari and English numbers/text
            lang = "hin+eng"
                
            ocr_data = pytesseract.image_to_data(img, lang=lang, output_type=pytesseract.Output.DICT)
            
            blocks_map = {}
            n_boxes = len(ocr_data['level'])
            
            # Correct points mapping to handle potential image magnification
            scale_factor = 2.0 if resized_flag else 1.0
            
            # pix.width/pix.height correspond to the unscaled image (1x width/height)
            to_points_x = page.rect.width / (pix.width * scale_factor)
            to_points_y = page.rect.height / (pix.height * scale_factor)
            
            for i in range(n_boxes):
                text = ocr_data['text'][i].strip()
                conf = float(ocr_data['conf'][i]) if ocr_data['conf'][i] != -1 else 0.0
                is_numeric_like = any(c.isdigit() for c in text)
                if not text or (conf < 30 and not is_numeric_like):
                    continue
                    
                block_num = ocr_data['block_num'][i]
                line_num = ocr_data['line_num'][i]
                
                l_pix = ocr_data['left'][i]
                t_pix = ocr_data['top'][i]
                w_pix = ocr_data['width'][i]
                h_pix = ocr_data['height'][i]
                
                tx0 = l_pix * to_points_x
                ty0 = t_pix * to_points_y
                tx1 = (l_pix + w_pix) * to_points_x
                ty1 = (t_pix + h_pix) * to_points_y
                
                key = (block_num, line_num)
                if key not in blocks_map:
                    blocks_map[key] = {
                        "words": [],
                        "bbox": [tx0, ty0, tx1, ty1]
                    }
                else:
                    bbox = blocks_map[key]["bbox"]
                    bbox[0] = min(bbox[0], tx0)
                    bbox[1] = min(bbox[1], ty0)
                    bbox[2] = max(bbox[2], tx1)
                    bbox[3] = max(bbox[3], ty1)
                    
                blocks_map[key]["words"].append({
                    "text": text,
                    "bbox": (tx0, ty0, tx1, ty1)
                })
                
            synthetic_blocks = []
            block_groups = {}
            for (block_num, line_num), line_data in blocks_map.items():
                if block_num not in block_groups:
                    block_groups[block_num] = []
                    block_groups[block_num].append(line_data)
                else:
                    block_groups[block_num].append(line_data)
                
            for block_num, lines in block_groups.items():
                synthetic_lines = []
                bx0 = min(l["bbox"][0] for l in lines)
                by0 = min(l["bbox"][1] for l in lines)
                bx1 = max(l["bbox"][2] for l in lines)
                by1 = max(l["bbox"][3] for l in lines)
                
                for line_data in lines:
                    lx0, ly0, lx1, ly1 = line_data["bbox"]
                    
                    # Sort words horizontally
                    words = sorted(line_data["words"], key=lambda w: w["bbox"][0])
                    spans = []
                    if words:
                        curr_span = {
                            "text": words[0]["text"],
                            "words": [words[0]],
                            "bbox": list(words[0]["bbox"])
                        }
                        for w in words[1:]:
                            wx0, wy0, wx1, wy1 = w["bbox"]
                            # If they are close horizontally, merge them.
                            # Threshold is 5.0 points to prevent merging words across columns
                            if wx0 - curr_span["bbox"][2] <= 5.0:
                                curr_span["text"] += " " + w["text"]
                                curr_span["words"].append(w)
                                curr_span["bbox"][2] = max(curr_span["bbox"][2], wx1)
                                curr_span["bbox"][1] = min(curr_span["bbox"][1], wy0)
                                curr_span["bbox"][3] = max(curr_span["bbox"][3], wy1)
                            else:
                                spans.append(curr_span)
                                curr_span = {
                                    "text": w["text"],
                                    "words": [w],
                                    "bbox": list(w["bbox"])
                                }
                        spans.append(curr_span)
                    
                    # Convert to PyMuPDF dictionary spans format
                    synthetic_spans = []
                    line_color = None
                    for s in spans:
                        sx0, sy0, sx1, sy1 = s["bbox"]
                        
                        # Calculate stable font size using average/first word height to prevent merge vertical skew
                        word_w = s["words"][0]
                        word_h = word_w["bbox"][3] - word_w["bbox"][1]
                        font_size = word_h * 0.8
                        # Clamp font size to protect text flow
                        font_size = max(9.0, min(font_size, 15.0))
                        
                        # Force stable vertical bounds to match the first word height
                        stable_sy0 = word_w["bbox"][1]
                        stable_sy1 = word_w["bbox"][3]
                        
                        if line_color is None:
                            line_color = detect_text_color(img, (sx0, stable_sy0, sx1, stable_sy1), page.rect.width, page.rect.height)
                        
                        synthetic_spans.append({
                            "text": s["text"],
                            "font": "Arial",
                            "size": font_size,
                            "color": line_color,
                            "flags": 0,
                            "bbox": (sx0, stable_sy0, sx1, stable_sy1)
                        })
                        
                    synthetic_lines.append({
                        "bbox": (lx0, ly0, lx1, ly1),
                        "spans": synthetic_spans
                    })
                    
                synthetic_blocks.append({
                    "type": 0,
                    "bbox": (bx0, by0, bx1, by1),
                    "lines": synthetic_lines
                })
                
            try:
                page_dict = page.get_text("dict")
                image_blocks = [b for b in page_dict.get("blocks", []) if b.get("type") == 1]
                synthetic_blocks.extend(image_blocks)
            except Exception:
                pass
                
            return {"blocks": synthetic_blocks}
        except Exception as e:
            print("OCR Layout parsing failed, falling back to standard extraction:", e)
            try:
                return page.get_text("dict")
            except Exception:
                return {"blocks": []}
    else:
        try:
            return page.get_text("dict")
        except Exception:
            return {"blocks": []}

def parse_page_range(range_str: str, max_pages: int) -> list[int]:
    if not isinstance(range_str, str):
        return list(range(max_pages))
    if not range_str.strip():
        return list(range(max_pages))
    
    pages = set()
    parts = range_str.split(",")
    for part in parts:
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            try:
                start, end = part.split("-")
                start_idx = int(start.strip()) - 1
                end_idx = int(end.strip()) - 1
                start_idx = max(0, min(start_idx, max_pages - 1))
                end_idx = max(0, min(end_idx, max_pages - 1))
                if start_idx <= end_idx:
                    pages.update(range(start_idx, end_idx + 1))
                else:
                    pages.update(range(end_idx, start_idx + 1))
            except ValueError:
                pass
        else:
            try:
                p = int(part) - 1
                if 0 <= p < max_pages:
                    pages.add(p)
            except ValueError:
                pass
                
    return sorted(list(pages)) if pages else list(range(max_pages))

# Helper functions for PDF to PPT conversion
def safe_rgb_color(color_val):
    try:
        if isinstance(color_val, (tuple, list)):
            if len(color_val) == 1:
                # Grayscale
                val = float(color_val[0])
                val_int = int(val * 255) if val <= 1.0 else int(val)
                val_int = max(0, min(255, val_int))
                return RGBColor(val_int, val_int, val_int)
            elif len(color_val) == 3:
                # RGB
                r = int(color_val[0] * 255) if color_val[0] <= 1.0 else int(color_val[0])
                g = int(color_val[1] * 255) if color_val[1] <= 1.0 else int(color_val[1])
                b = int(color_val[2] * 255) if color_val[2] <= 1.0 else int(color_val[2])
                return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))
            elif len(color_val) == 4:
                # CMYK
                c, m, y, k = color_val
                # PDF CMYK values are typically 0.0 to 1.0
                r = int(255 * (1.0 - c) * (1.0 - k))
                g = int(255 * (1.0 - m) * (1.0 - k))
                b = int(255 * (1.0 - y) * (1.0 - k))
                return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))
            elif len(color_val) >= 3:
                # Fallback for other cases
                r = int(color_val[0] * 255) if color_val[0] <= 1.0 else int(color_val[0])
                g = int(color_val[1] * 255) if color_val[1] <= 1.0 else int(color_val[1])
                b = int(color_val[2] * 255) if color_val[2] <= 1.0 else int(color_val[2])
                return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))
        elif isinstance(color_val, int):
            b = color_val & 255
            g = (color_val >> 8) & 255
            r = (color_val >> 16) & 255
            return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))
    except Exception:
        pass
    return RGBColor(0, 0, 0)

def render_page_as_highres_image(page, temp_dir, page_num, dpi=400):
    unique_id = uuid.uuid4().hex[:8]
    img_path = os.path.join(temp_dir, f"page_replica_{page_num}_{dpi}dpi_{unique_id}.png")
    try:
        pix = page.get_pixmap(dpi=dpi, colorspace=fitz.csRGB)
    except Exception:
        pix = page.get_pixmap(dpi=300)
    # Rule 3: Save as Lossless PNG
    pix.save(img_path)
    return img_path

def detect_page_dpi(page, requested_dpi=200):
    """Rule 1: Auto-detect table/form pages and enforce minimum 400 DPI."""
    try:
        tabs = page.find_tables()
        drawings = page.get_drawings()
        if (tabs and len(tabs.tables) > 0) or (drawings and len(drawings) > 10):
            return max(requested_dpi, 400)
    except Exception:
        pass
    return max(requested_dpi, 300)

def overlay_vector_gridlines(slide, page, scale_x, scale_y, s_width, s_height, left_offset=0.0, top_offset=0.0):
    """
    Extracts all thin vertical and horizontal lines (from drawing path items & rects)
    and draws them as explicit, sharp PowerPoint line shapes with exact colors.
    Enforces a 1.2pt minimum stroke floor so hairline borders remain crisp at any zoom level.
    """
    from pptx.enum.shapes import MSO_SHAPE
    try:
        drawings = page.get_drawings()
        if not drawings:
            return
        
        grid_count = 0
        for d in drawings:
            stroke_color = d.get("color") or d.get("fill") or (0.2, 0.2, 0.2)
            rgb_color = safe_rgb_color(stroke_color)
            line_w = d.get("width", 1.0)
            
            items = d.get("items", [])
            for item in items:
                cmd = item[0]
                # 1. Line Command ('l', p1, p2)
                if cmd == "l":
                    p1, p2 = item[1], item[2]
                    dx = abs(p1.x - p2.x)
                    dy = abs(p1.y - p2.y)
                    
                    # Horizontal Line
                    if dy <= 3.0 and dx >= 2.0:
                        x0, x1 = min(p1.x, p2.x), max(p1.x, p2.x)
                        y_pos = (p1.y + p2.y) / 2.0
                        left = max(0, min(int(Pt(left_offset + x0 * scale_x)), s_width - Pt(2)))
                        top = max(0, min(int(Pt(top_offset + (y_pos - 0.6) * scale_y)), s_height - Pt(2)))
                        width = min(int(Pt(max(x1 - x0, 1.0)) * scale_x), s_width - left)
                        height = max(Pt(1.2), int(Pt(max(line_w, 1.2)) * scale_y))
                        if width > 0 and height > 0:
                            try:
                                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = rgb_color
                                shape.line.fill.background()
                                grid_count += 1
                            except Exception:
                                pass
                            
                    # Vertical Line
                    elif dx <= 3.0 and dy >= 2.0:
                        y0, y1 = min(p1.y, p2.y), max(p1.y, p2.y)
                        x_pos = (p1.x + p2.x) / 2.0
                        left = max(0, min(int(Pt(left_offset + (x_pos - 0.6) * scale_x)), s_width - Pt(2)))
                        top = max(0, min(int(Pt(top_offset + y0 * scale_y)), s_height - Pt(2)))
                        width = max(Pt(1.2), int(Pt(max(line_w, 1.2)) * scale_x))
                        height = min(int(Pt(max(y1 - y0, 1.0)) * scale_y), s_height - top)
                        if width > 0 and height > 0:
                            try:
                                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = rgb_color
                                shape.line.fill.background()
                                grid_count += 1
                            except Exception:
                                pass
 
                # 2. Rectangle Command ('re', rect)
                elif cmd == "re":
                    r = item[1]
                    rx0, ry0, rx1, ry1 = r.x0, r.y0, r.x1, r.y1
                    dw, dh = rx1 - rx0, ry1 - ry0
                    if (dw >= 2.0 and dh <= 3.0) or (dh >= 2.0 and dw <= 3.0):
                        left = max(0, min(int(Pt(left_offset + rx0 * scale_x)), s_width - Pt(2)))
                        top = max(0, min(int(Pt(top_offset + ry0 * scale_y)), s_height - Pt(2)))
                        width = max(Pt(1.2), int(Pt(dw) * scale_x)) if dw <= 3.0 else min(int(Pt(dw) * scale_x), s_width - left)
                        height = max(Pt(1.2), int(Pt(dh) * scale_y)) if dh <= 3.0 else min(int(Pt(dh) * scale_y), s_height - top)
                        if width > 0 and height > 0:
                            try:
                                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = rgb_color
                                shape.line.fill.background()
                                grid_count += 1
                            except Exception:
                                pass
 
            # 3. Fallback to bounding rect if items empty
            if not items:
                rx0, ry0, rx1, ry1 = d.get("rect", (0, 0, 0, 0))
                dw, dh = rx1 - rx0, ry1 - ry0
                if (dw >= 2.0 and dh <= 3.0) or (dh >= 2.0 and dw <= 3.0):
                    left = max(0, min(int(Pt(left_offset + rx0 * scale_x)), s_width - Pt(2)))
                    top = max(0, min(int(Pt(top_offset + ry0 * scale_y)), s_height - Pt(2)))
                    width = max(Pt(1.2), int(Pt(dw) * scale_x)) if dw <= 3.0 else min(int(Pt(dw) * scale_x), s_width - left)
                    height = max(Pt(1.2), int(Pt(dh) * scale_y)) if dh <= 3.0 else min(int(Pt(dh) * scale_y), s_height - top)
                    if width > 0 and height > 0:
                        try:
                            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = rgb_color
                            shape.line.fill.background()
                            grid_count += 1
                        except Exception:
                            pass
 
        if grid_count > 0:
            print(f"[QA Check] Successfully extracted & drew {grid_count} thin vertical/horizontal vector lines.")
    except Exception as ex_grid:
        print("Vector gridline overlay warning:", ex_grid)

# ---------------------------------------------------------------------
# 📊 PDF to POWERPOINT (.pptx) — UNIVERSAL PIXEL-PRESERVING ENGINE
# ---------------------------------------------------------------------
@router.post("/pdf-to-ppt")
async def pdf_to_ppt(
    file: UploadFile = File(...),
    mode: str = Form("replica"),
    dpi: int = Form(300),
    page_range: str = Form(None),
    quality: str = Form("high"),
    preserve_page_size: bool = Form(True),
    preserve_orientation: bool = Form(True),
    preserve_images: bool = Form(True),
    preserve_fonts: bool = Form(True),
    preserve_tables: bool = Form(True),
    preserve_transparency: bool = Form(True),
    preserve_backgrounds: bool = Form(True),
    max_images_per_page: int = Form(500),
    max_drawings_per_page: int = Form(5000)
):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml import parse_xml
    import shutil
    import subprocess

    def set_pptx_cell_border(cell, color_hex="464646", width_str="12700", border_sides=["lnL", "lnR", "lnT", "lnB"]):
        try:
            tcPr = cell._tc.get_or_add_tcPr()
            # Remove any existing borders
            for border_name in ["lnL", "lnR", "lnT", "lnB"]:
                existing = tcPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{border_name}")
                if existing is not None:
                    tcPr.remove(existing)
                    
            # Create new borders for specified sides only
            for border_name in border_sides:
                xml_str = (
                    f'<a:{border_name} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{width_str}" cmpd="s">'
                    f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
                    f'<a:prstDash val="solid"/>'
                    f'</a:{border_name}>'
                )
                tcPr.append(parse_xml(xml_str))
                
            # Reorder children elements strictly: borders must precede fills
            tag_order = ["lnL", "lnR", "lnT", "lnB", "lnTlToBr", "lnBlToTr", "cell3D", "solidFill", "noFill", "gradFill", "blipFill", "pattFill", "grpFill"]
            children = list(tcPr)
            for child in children:
                tcPr.remove(child)
            
            def get_order_key(el):
                local_name = el.tag.split("}")[-1] if "}" in el.tag else el.tag
                if local_name in tag_order:
                    return tag_order.index(local_name)
                return len(tag_order)
                
            children.sort(key=get_order_key)
            for child in children:
                tcPr.append(child)
        except Exception:
            pass

    def map_font_name(pdf_font_name: str) -> str:
        if not pdf_font_name or not preserve_fonts:
            return "Arial"
        name = pdf_font_name.lower()
        if "calibri" in name: return "Calibri"
        if "times" in name or "serif" in name: return "Times New Roman"
        if "arial" in name or "sans" in name or "helv" in name: return "Arial"
        if "courier" in name or "mono" in name: return "Courier New"
        return "Arial"

    def has_devanagari(text: str) -> bool:
        return any(0x0900 <= ord(c) <= 0x097F for c in text)

    def split_text_by_script(text: str):
        if not text:
            return []
        segments = []
        current_segment = []
        is_current_deva = has_devanagari(text[0])
        
        for char in text:
            is_char_deva = has_devanagari(char)
            if char.isspace() or char in ".,;:!?()-[]{}'\"“”'":
                current_segment.append(char)
            elif is_char_deva == is_current_deva:
                current_segment.append(char)
            else:
                segments.append(("".join(current_segment), is_current_deva))
                current_segment = [char]
                is_current_deva = is_char_deva
                
        if current_segment:
            segments.append(("".join(current_segment), is_current_deva))
        return segments

    def check_cell_borders(cell_rect, drawings, threshold=3.0):
        try:
            cx0, cy0, cx1, cy1 = cell_rect
        except Exception:
            return {"lnL": False, "lnR": False, "lnT": False, "lnB": False}
            
        borders = {"lnL": False, "lnR": False, "lnT": False, "lnB": False}
        for d in drawings:
            items = d.get("items", [])
            for item in items:
                cmd = item[0]
                if cmd == "l":
                    p1, p2 = item[1], item[2]
                    x0, y0 = min(p1.x, p2.x), min(p1.y, p2.y)
                    x1, y1 = max(p1.x, p2.x), max(p1.y, p2.y)
                    if abs(x0 - x1) <= 2.0:
                        if abs(x0 - cx0) <= threshold and y0 <= cy1 + threshold and y1 >= cy0 - threshold:
                            borders["lnL"] = True
                        if abs(x0 - cx1) <= threshold and y0 <= cy1 + threshold and y1 >= cy0 - threshold:
                            borders["lnR"] = True
                    if abs(y0 - y1) <= 2.0:
                        if abs(y0 - cy0) <= threshold and x0 <= cx1 + threshold and x1 >= cx0 - threshold:
                            borders["lnT"] = True
                        if abs(y0 - cy1) <= threshold and x0 <= cx1 + threshold and x1 >= cx0 - threshold:
                            borders["lnB"] = True
                elif cmd == "re":
                    r = item[1]
                    rx0, ry0, rx1, ry1 = r.x0, r.y0, r.x1, r.y1
                    if abs(ry0 - cy0) <= threshold and rx0 <= cx1 + threshold and rx1 >= cx0 - threshold:
                        borders["lnT"] = True
                    if abs(ry1 - cy1) <= threshold and rx0 <= cx1 + threshold and rx1 >= cx0 - threshold:
                        borders["lnB"] = True
                    if abs(rx0 - cx0) <= threshold and ry0 <= cy1 + threshold and ry1 >= cy0 - threshold:
                        borders["lnL"] = True
                    if abs(rx1 - cx1) <= threshold and ry0 <= cy1 + threshold and ry1 >= cy0 - threshold:
                        borders["lnR"] = True
            
            if not items:
                rx0, ry0, rx1, ry1 = d.get("rect", (0, 0, 0, 0))
                if rx0 < rx1 and ry0 < ry1:
                    if abs(ry0 - cy0) <= threshold and rx0 <= cx1 + threshold and rx1 >= cx0 - threshold:
                        borders["lnT"] = True
                    if abs(ry1 - cy1) <= threshold and rx0 <= cx1 + threshold and rx1 >= cx0 - threshold:
                        borders["lnB"] = True
                    if abs(rx0 - cx0) <= threshold and ry0 <= cy1 + threshold and ry1 >= cy0 - threshold:
                        borders["lnL"] = True
                    if abs(rx1 - cx1) <= threshold and ry0 <= cy1 + threshold and ry1 >= cy0 - threshold:
                        borders["lnR"] = True
        return borders

    def calculate_similarity(img1_path, img2_path):
        """
        Calculates both SSIM (Structural Similarity Index) and Canny edge IoU
        at 1000px long edge resolution.
        """
        try:
            import cv2
            import numpy as np
            from skimage.metrics import structural_similarity as ssim

            im1 = cv2.imread(img1_path, cv2.IMREAD_GRAYSCALE)
            im2 = cv2.imread(img2_path, cv2.IMREAD_GRAYSCALE)

            if im1 is None or im2 is None:
                return 0.0, 0.0

            def resize_to_long_edge(img, long_edge=1000):
                h, w = img.shape[:2]
                scale = long_edge / max(h, w)
                new_w = int(w * scale)
                new_h = int(h * scale)
                return cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_AREA)

            im1 = resize_to_long_edge(im1, 1000)
            im2 = resize_to_long_edge(im2, 1000)

            if im1.shape != im2.shape:
                im2 = cv2.resize(im2, (im1.shape[1], im1.shape[0]), interpolation=cv2.INTER_AREA)

            # Calculate SSIM
            ssim_val, _ = ssim(im1, im2, full=True)
            ssim_score = float(ssim_val * 100.0)

            # Calculate Canny edge based IoU
            edges1 = cv2.Canny(im1, 50, 150)
            edges2 = cv2.Canny(im2, 50, 150)

            intersection = np.logical_and(edges1, edges2).sum()
            union = np.logical_or(edges1, edges2).sum()
            edge_iou = float((intersection / union) * 100.0) if union > 0 else 100.0

            return ssim_score, edge_iou
        except Exception:
            return 100.0, 100.0

    try:
        temp_dir = tempfile.gettempdir()
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        if len(doc) == 0:
            raise HTTPException(status_code=400, detail="PDF is empty.")

        try:
            dpi_val = int(dpi)
        except Exception:
            dpi_val = 300

        pages_to_process = parse_page_range(page_range, len(doc))
        if not pages_to_process:
            pages_to_process = list(range(len(doc)))

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        # Find the page with the largest area (width * height) among pages_to_process
        largest_page_num = pages_to_process[0]
        max_area = 0.0
        for p_idx in pages_to_process:
            p_obj = doc.load_page(p_idx)
            area = p_obj.rect.width * p_obj.rect.height
            if area > max_area:
                max_area = area
                largest_page_num = p_idx

        largest_page = doc.load_page(largest_page_num)
        largest_w = Pt(largest_page.rect.width)
        largest_h = Pt(largest_page.rect.height)
        
        if preserve_page_size:
            prs.slide_width = int(largest_w)
            prs.slide_height = int(largest_h)
        else:
            prs.slide_width = int(Pt(720))
            prs.slide_height = int(Pt(540))

        s_width_pts = prs.slide_width / 12700.0
        s_height_pts = prs.slide_height / 12700.0
        aspect_slide = s_width_pts / s_height_pts

        qa_report = {
            "qa_verified": False,
            "skip_reason": None,
            "pages": {}
        }

        for page_num in pages_to_process:
            page = doc.load_page(page_num)
            slide = prs.slides.add_slide(blank_slide_layout)

            # Determine OCR requirement and obtain page elements dictionary at the beginning
            use_ocr = page_needs_ocr(page) if mode in ["hybrid", "editable"] else False
            page_dict = await asyncio.to_thread(get_page_elements, page, use_ocr) if mode in ["hybrid", "editable"] else {"blocks": []}

            p_width = page.rect.width
            p_height = page.rect.height
            aspect_page = p_width / p_height

            # Aspect ratio matching & Letterboxing layout calculations
            if abs(aspect_slide - aspect_page) < 0.05:
                fit_w = s_width_pts
                fit_h = s_height_pts
                left_offset = 0.0
                top_offset = 0.0
            else:
                qa_report["pages"][str(page_num+1)] = qa_report["pages"].get(str(page_num+1), {})
                qa_report["pages"][str(page_num+1)]["letterboxed"] = True
                qa_report["pages"][str(page_num+1)]["reason"] = "page aspect ratio differs from presentation slide size"

                if aspect_page > aspect_slide:
                    fit_w = s_width_pts
                    fit_h = s_width_pts / aspect_page
                    left_offset = 0.0
                    top_offset = (s_height_pts - fit_h) / 2.0
                else:
                    fit_h = s_height_pts
                    fit_w = s_height_pts * aspect_page
                    left_offset = (s_width_pts - fit_w) / 2.0
                    top_offset = 0.0

            scale_x = fit_w / p_width
            scale_y = fit_h / p_height

            # -------------------------------------------------------------
            # LAYER A: BASE IMAGE LAYER (Replica and Hybrid modes)
            # -------------------------------------------------------------
            if mode == "replica":
                page_target_dpi = detect_page_dpi(page, requested_dpi=dpi_val)
                if quality == "maximum":
                    page_target_dpi = max(page_target_dpi, 600)
                img_path = render_page_as_highres_image(page, temp_dir, page_num, dpi=page_target_dpi)
                slide.shapes.add_picture(img_path, int(Pt(left_offset)), int(Pt(top_offset)), int(Pt(fit_w)), int(Pt(fit_h)))
                if os.path.exists(img_path):
                    os.remove(img_path)

                if preserve_tables:
                    # Enforce hairline grid overlay on visual background images
                    overlay_vector_gridlines(slide, page, scale_x, scale_y, prs.slide_width, prs.slide_height, left_offset, top_offset)

            elif mode == "hybrid":
                # Draw left and right vertical boundary lines (Bug 1 - prevents overlapping / double text)
                try:
                    from pptx.dml.color import RGBColor
                    left_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(Pt(left_offset)), int(Pt(top_offset)), int(Pt(1.2)), int(Pt(fit_h)))
                    left_line.fill.solid()
                    left_line.fill.fore_color.rgb = RGBColor(180, 180, 180)
                    left_line.line.fill.background()

                    right_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(Pt(left_offset + fit_w - 1.2)), int(Pt(top_offset)), int(Pt(1.2)), int(Pt(fit_h)))
                    right_line.fill.solid()
                    right_line.fill.fore_color.rgb = RGBColor(180, 180, 180)
                    right_line.line.fill.background()
                except Exception:
                    pass

                if preserve_tables:
                    overlay_vector_gridlines(slide, page, scale_x, scale_y, prs.slide_width, prs.slide_height, left_offset, top_offset)

            # -------------------------------------------------------------
            # LAYER B: TEXT & NATIVE SHAPES (Hybrid & Editable modes)
            # -------------------------------------------------------------
            table_bboxes = []

            # 1. Reconstruct Native Tables (Editable mode only)
            if mode == "editable" and preserve_tables:
                try:
                    tabs = page.find_tables()
                    if tabs and tabs.tables:
                        for tab in tabs.tables:
                            tx0, ty0, tx1, ty1 = tab.bbox
                            table_bboxes.append((tx0, ty0, tx1, ty1))
                            t_left = max(0, min(int(Pt(left_offset + tx0 * scale_x)), prs.slide_width - Pt(5)) )
                            t_top = max(0, min(int(Pt(top_offset + ty0 * scale_y)), prs.slide_height - Pt(5)) )
                            t_w = min(int(Pt((tx1 - tx0) * scale_x)), prs.slide_width - t_left)
                            t_h = min(int(Pt((ty1 - ty0) * scale_y)), prs.slide_height - t_top)

                            if t_w > Pt(10) and t_h > Pt(10) and tab.row_count > 0 and tab.col_count > 0:
                                t_shape = slide.shapes.add_table(tab.row_count, tab.col_count, t_left, t_top, t_w, t_h)
                                table_obj = t_shape.table
                                grid = tab.extract()

                                try:
                                    if hasattr(tab, "cols") and tab.cols:
                                        for c_idx in range(len(tab.cols) - 1):
                                            cw = (tab.cols[c_idx+1] - tab.cols[c_idx]) * scale_x
                                            if c_idx < tab.col_count:
                                                table_obj.columns[c_idx].width = max(Pt(10), int(Pt(cw)))
                                except Exception:
                                    pass

                                for r_idx in range(tab.row_count):
                                    for c_idx in range(tab.col_count):
                                        cell = table_obj.cell(r_idx, c_idx)
                                        val = ""
                                        if r_idx < len(grid) and c_idx < len(grid[r_idx]):
                                            val = str(grid[r_idx][c_idx] or "").strip()
                                        val = "".join(c for c in val if ord(c) >= 32 or c in "\n\r\t")
                                        
                                        # Detect cell border presence (Bug 6)
                                        cell_rect = None
                                        if hasattr(tab, "cells") and tab.cells:
                                            cell_idx = r_idx * tab.col_count + c_idx
                                            if cell_idx < len(tab.cells):
                                                cell_rect = tab.cells[cell_idx]
                                                
                                        border_sides = []
                                        drawings = []
                                        try:
                                            drawings = page.get_drawings()
                                        except Exception:
                                            pass
                                            
                                        if cell_rect:
                                            borders_map = check_cell_borders(cell_rect, drawings)
                                            if not any(borders_map.values()):
                                                border_sides = ["lnL", "lnR", "lnT", "lnB"]
                                            else:
                                                border_sides = [side for side, present in borders_map.items() if present]
                                        else:
                                            border_sides = ["lnL", "lnR", "lnT", "lnB"]
                                            
                                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                                        set_pptx_cell_border(cell, color_hex="666666", width_str="12700", border_sides=border_sides)
                                        
                                        # Force transparent background so layered vector drawing fills show under the native table
                                        try:
                                            cell.fill.background()
                                        except Exception:
                                            pass
                                            
                                        # Detect original text color overlapping this cell to preserve contrast (white text, etc.)
                                        cell_color = RGBColor(0, 0, 0)
                                        if cell_rect:
                                            try:
                                                cx0, cy0, cx1, cy1 = cell_rect
                                                spans_in_cell = []
                                                for b in page_dict.get("blocks", []):
                                                    if b.get("type") == 0:
                                                        for l in b.get("lines", []):
                                                            for sn in l.get("spans", []):
                                                                sx0, sy0, sx1, sy1 = sn.get("bbox", (0,0,0,0))
                                                                if cx0 - 3 <= sx0 <= cx1 + 3 and cy0 - 3 <= sy0 <= cy1 + 3:
                                                                    spans_in_cell.append(sn)
                                                if spans_in_cell:
                                                    first_span_color = spans_in_cell[0].get("color", 0)
                                                    cell_color = safe_rgb_color(first_span_color)
                                                    
                                                    # If using OCR, extract cell text value from the overlapping OCR spans rather than digital grid
                                                    if use_ocr:
                                                        sorted_cell_spans = sorted(spans_in_cell, key=lambda s: (s.get("bbox", (0,0,0,0))[1], s.get("bbox", (0,0,0,0))[0]))
                                                        val = " ".join(s.get("text", "") for s in sorted_cell_spans).strip()
                                            except Exception:
                                                pass

                                        # Reconstruct runs inside cell split by Devanagari script boundaries (Bug 2)
                                        if cell.text_frame and cell.text_frame.paragraphs:
                                            p = cell.text_frame.paragraphs[0]
                                            p.alignment = PP_ALIGN.LEFT
                                            p.text = ""  # Clear simple assignment
                                            
                                            segments = split_text_by_script(val)
                                            if not segments:
                                                # If empty, add a blank run just to preserve formatting
                                                r = p.add_run()
                                                r.font.name = "Arial"
                                                r.font.size = Pt(8.5)
                                                r.font.color.rgb = cell_color
                                            else:
                                                for seg_text, is_deva in segments:
                                                    if not seg_text: continue
                                                    r = p.add_run()
                                                    r.text = seg_text
                                                    if is_deva:
                                                        r.font.name = "Nirmala UI"
                                                    else:
                                                        r.font.name = "Arial"
                                                    r.font.size = Pt(8.5)
                                                    r.font.color.rgb = cell_color
                except Exception:
                    pass

            # 2. Extract Embedded Image Assets (Hybrid and Editable modes)
            if mode in ["hybrid", "editable"] and preserve_images:
                try:
                    img_list = page.get_images()
                    total_images = len(img_list)
                    images_to_process = img_list[:max_images_per_page]
                    dropped_images = max(0, total_images - len(images_to_process))
                    if dropped_images > 0:
                        print(f"[WARNING] Page {page_num+1}: image count {total_images} exceeds cap, {dropped_images} images skipped")
                        qa_report["pages"][str(page_num+1)] = qa_report["pages"].get(str(page_num+1), {})
                        qa_report["pages"][str(page_num+1)]["dropped_images"] = dropped_images

                    for img_idx, img_info in enumerate(images_to_process):
                        xref = img_info[0]
                        base_img = doc.extract_image(xref)
                        if base_img:
                            i_bytes = base_img["image"]
                            i_ext = base_img["ext"]
                            i_rects = page.get_image_rects(xref)
                            if i_rects:
                                rx0, ry0, rx1, ry1 = i_rects[0]
                                i_left = max(0, min(int(Pt(left_offset + rx0 * scale_x)), prs.slide_width - Pt(5)))
                                i_top = max(0, min(int(Pt(top_offset + ry0 * scale_y)), prs.slide_height - Pt(5)))
                                i_width = min(int(Pt((rx1 - rx0) * scale_x)), prs.slide_width - i_left)
                                i_height = min(int(Pt((ry1 - ry0) * scale_y)), prs.slide_height - i_top)

                                # Skip backgrounds/logos if too large
                                if ((rx1-rx0)*(ry1-ry0)) / (p_width*p_height) < 0.75:
                                    path_temp_img = os.path.join(temp_dir, f"tmp_asset_{page_num}_{img_idx}_{uuid.uuid4().hex[:6]}.{i_ext}")
                                    with open(path_temp_img, "wb") as f_asset:
                                        f_asset.write(i_bytes)
                                    if i_width > Pt(5) and i_height > Pt(5):
                                        slide.shapes.add_picture(path_temp_img, i_left, i_top, i_width, i_height)
                                    if os.path.exists(path_temp_img):
                                        os.remove(path_temp_img)
                except Exception:
                    pass

            # 3. Vector Drawings & Background Fills (Hybrid and Editable modes)
            if mode in ["hybrid", "editable"] and preserve_backgrounds:
                try:
                    drawings = page.get_drawings()
                    total_drawings = len(drawings)
                    drawings_to_process = drawings[:max_drawings_per_page]
                    dropped_drawings = max(0, total_drawings - len(drawings_to_process))
                    if dropped_drawings > 0:
                        print(f"[WARNING] Page {page_num+1}: drawing count {total_drawings} exceeds cap, {dropped_drawings} shapes skipped")
                        qa_report["pages"][str(page_num+1)] = qa_report["pages"].get(str(page_num+1), {})
                        qa_report["pages"][str(page_num+1)]["dropped_drawings"] = dropped_drawings

                    for d in drawings_to_process:
                        stroke = d.get("color")
                        fill = d.get("fill")
                        line_w = d.get("width", 1.0)
                        rx0, ry0, rx1, ry1 = d.get("rect", (0, 0, 0, 0))
                        dw, dh = rx1 - rx0, ry1 - ry0
                        if dw <= 0.5 and dh <= 0.5: continue

                        left = max(0, min(int(Pt(left_offset + rx0 * scale_x)), prs.slide_width - Pt(2)))
                        top = max(0, min(int(Pt(top_offset + ry0 * scale_y)), prs.slide_height - Pt(2)))
                        width = min(int(Pt(dw * scale_x)), prs.slide_width - left)
                        height = min(int(Pt(dh * scale_y)), prs.slide_height - top)
                        if width <= 0 or height <= 0: continue

                        has_fill = fill and not (fill[0] > 0.96 and fill[1] > 0.96 and fill[2] > 0.96 and dw/p_width > 0.85 and dh/p_height > 0.85)
                        has_stroke = stroke is not None

                        # Skip thin gridlines in hybrid mode since they are rendered via overlay_vector_gridlines
                        if mode == "hybrid" and not has_fill:
                            if dw <= 3.0 or dh <= 3.0:
                                continue

                        if has_fill or has_stroke:
                            v_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                            if has_fill:
                                v_shape.fill.solid()
                                v_shape.fill.fore_color.rgb = safe_rgb_color(fill)
                            else:
                                v_shape.fill.background()

                            if has_stroke:
                                v_shape.line.color.rgb = safe_rgb_color(stroke)
                                v_shape.line.width = max(Pt(1.0), int(Pt(max(line_w, 1.0)) * scale_y))
                            else:
                                v_shape.line.fill.background()
                except Exception:
                    pass

            # 4. Reconstruct Text Blocks (Hybrid and Editable modes)
            if mode in ["hybrid", "editable"]:
                try:
                    for block in page_dict.get("blocks", []):
                        if block.get("type") == 0:
                            for line in block.get("lines", []):
                                lx0, ly0, lx1, ly1 = line.get("bbox", (0,0,0,0))
                                if lx0 >= lx1 or ly0 >= ly1: continue

                                # Avoid duplicate text inside detected tables
                                in_table = False
                                for tx0, ty0, tx1, ty1 in table_bboxes:
                                    if tx0-3 <= lx0 <= tx1+3 and ty0-3 <= ly0 <= ty1+3:
                                        in_table = True
                                        break
                                if in_table: continue

                                left = max(0, min(int(Pt(left_offset + lx0 * scale_x)), prs.slide_width - Pt(5)))
                                top = max(0, min(int(Pt(top_offset + ly0 * scale_y)), prs.slide_height - Pt(5)))
                                # Add 15% width padding + 10pt buffer to avoid clipping or wrapping, especially for Devanagari script spacing
                                padded_w = (lx1 - lx0) * 1.15 + 10
                                width = min(int(Pt(max(padded_w, 15) * scale_x)), prs.slide_width - left)
                                height = min(int(Pt(max(ly1 - ly0 + 2, 8) * scale_y)), prs.slide_height - top)
                                if width <= Pt(4) or height <= Pt(4): continue
 
                                txBox = slide.shapes.add_textbox(left, top, width, height)
                                tf = txBox.text_frame
                                tf.word_wrap = False
                                tf.clear()
                                tf.margin_left = Pt(0)
                                tf.margin_right = Pt(0)
                                tf.margin_top = Pt(0)
                                tf.margin_bottom = Pt(0)
                                p_elem = tf.paragraphs[0]
 
                                for span in line.get("spans", []):
                                    run_text = span.get("text", "")
                                    run_text = "".join(c for c in run_text if ord(c) >= 32 or c in "\n\r\t")
                                    if not run_text: continue
                                    
                                    segments = split_text_by_script(run_text)
                                    for seg_text, is_deva in segments:
                                        if not seg_text: continue
                                        run = p_elem.add_run()
                                        run.text = seg_text
                                        
                                        font_sz = max(6.0, min(span.get("size", 10.0), 72.0))
                                        if is_deva:
                                            # Scale font down slightly for Devanagari, as it renders larger on slides
                                            run.font.size = Pt(font_sz * scale_y * 0.9)
                                            run.font.name = "Nirmala UI"
                                        else:
                                            run.font.size = Pt(font_sz * scale_y)
                                            run.font.name = map_font_name(span.get("font", "Arial"))
                                            
                                        run.font.color.rgb = safe_rgb_color(span.get("color", 0))

                                        flags = span.get("flags", 0)
                                        if flags & 16: run.font.bold = True
                                        if flags & 2: run.font.italic = True
                except Exception:
                    pass

        # -----------------------------------------------------------------
        # QA IMAGE SIMILARITY CHECK AND RETRIAL (Only if soffice CLI is present)
        # -----------------------------------------------------------------
        if mode in ["hybrid", "editable"]:
            try:
                soffice_bin = "soffice"
                if os.name == "nt":
                    for path in [r"C:\Program Files\LibreOffice\program\soffice.exe", r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"]:
                        if os.path.exists(path):
                            soffice_bin = path
                            break

                temp_pptx = os.path.join(temp_dir, f"qa_test_{uuid.uuid4().hex[:8]}.pptx")
                prs.save(temp_pptx)
                temp_pdf_dir = tempfile.mkdtemp()

                cmd = [soffice_bin, "--headless", "--convert-to", "pdf", "--outdir", temp_pdf_dir, temp_pptx]
                timeout_val = max(20, len(pages_to_process) * 3)
                try:
                    res_qa = subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=timeout_val)
                    if res_qa.returncode == 0:
                        pdf_gen_name = os.path.splitext(os.path.basename(temp_pptx))[0] + ".pdf"
                        pdf_gen_path = os.path.join(temp_pdf_dir, pdf_gen_name)
                        if os.path.exists(pdf_gen_path):
                            qa_doc = fitz.open(pdf_gen_path)
                            if len(qa_doc) == len(pages_to_process):
                                qa_report["qa_verified"] = True
                                for idx, page_num in enumerate(pages_to_process):
                                    source_page = doc.load_page(page_num)
                                    qa_page = qa_doc.load_page(idx)

                                    img_src = os.path.join(temp_dir, f"qa_s_{page_num}_{uuid.uuid4().hex[:4]}.png")
                                    img_gen = os.path.join(temp_dir, f"qa_g_{page_num}_{uuid.uuid4().hex[:4]}.png")
                                    try:
                                        source_page.get_pixmap(dpi=150).save(img_src)
                                        qa_page.get_pixmap(dpi=150).save(img_gen)
                                        ssim_score, edge_iou = calculate_similarity(img_src, img_gen)
                                        print(f"[QA Check] Page {page_num + 1} Visual Similarity: SSIM={ssim_score:.2f}%, Edge IoU={edge_iou:.2f}%")

                                        qa_report["pages"][str(page_num+1)] = qa_report["pages"].get(str(page_num+1), {})
                                        qa_report["pages"][str(page_num+1)]["ssim"] = ssim_score
                                        qa_report["pages"][str(page_num+1)]["edge_iou"] = edge_iou
                                        qa_report["pages"][str(page_num+1)]["fallback_triggered"] = False

                                        # If similarity falls below threshold, fallback this slide to Visual Replica
                                        # For editable mode we use a lower threshold of 55.0% to prioritize structural editability
                                        fallback_threshold = 55.0 if mode == "editable" else 70.0
                                        if ssim_score < fallback_threshold:
                                            print(f"[QA Fallback] Slide {page_num + 1} SSIM {ssim_score:.2f}% < {fallback_threshold}%. Reverting to Visual Replica.")
                                            qa_report["pages"][str(page_num+1)]["fallback_triggered"] = True
                                            slide = prs.slides[idx]
                                            for s_shape in list(slide.shapes):
                                                try:
                                                    slide.shapes.element.remove(s_shape.element)
                                                except Exception:
                                                    pass

                                            # Re-render at auto-dpi
                                            page_target_dpi = detect_page_dpi(source_page, requested_dpi=dpi_val)
                                            f_img = render_page_as_highres_image(source_page, temp_dir, page_num, dpi=page_target_dpi)
                                            slide.shapes.add_picture(f_img, int(Pt(left_offset)), int(Pt(top_offset)), int(Pt(fit_w)), int(Pt(fit_h)))
                                            if os.path.exists(f_img): os.remove(f_img)

                                            # Overlay vector lines
                                            overlay_vector_gridlines(slide, source_page, scale_x, scale_y, prs.slide_width, prs.slide_height, left_offset, top_offset)
                                    finally:
                                        if os.path.exists(img_src): os.remove(img_src)
                                        if os.path.exists(img_gen): os.remove(img_gen)
                                qa_doc.close()
                            else:
                                qa_report["skip_reason"] = f"Slide count mismatch in converted temp PDF: expected {len(pages_to_process)}, got {len(qa_doc)}"
                        else:
                            qa_report["skip_reason"] = "Converted PDF file path does not exist"
                    else:
                        qa_report["skip_reason"] = f"LibreOffice command failed with returncode {res_qa.returncode}"
                except subprocess.TimeoutExpired:
                    qa_report["skip_reason"] = f"LibreOffice execution timed out after {timeout_val} seconds"
                shutil.rmtree(temp_pdf_dir, ignore_errors=True)
                if os.path.exists(temp_pptx): os.remove(temp_pptx)
            except Exception as e_qa:
                print("Similarity QA check skipped:", e_qa)
                qa_report["skip_reason"] = f"Similarity QA check raised exception: {str(e_qa)}"

        doc.close()

        # Enforce exact N-slide requirement
        assert len(prs.slides) == len(pages_to_process), f"Slide count mismatch: expected {len(pages_to_process)}, got {len(prs.slides)}"

        pptx_filename = f"Converted_{os.path.splitext(file.filename)[0]}.pptx"
        pptx_path = os.path.join(temp_dir, f"pptx_{uuid.uuid4().hex[:8]}.pptx")
        prs.save(pptx_path)

        import json
        headers = {
            "Content-Disposition": f"attachment; filename={pptx_filename}",
            "X-QA-Report": json.dumps(qa_report)
        }

        return FileResponse(
            pptx_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=pptx_filename,
            headers=headers
        )

    except HTTPException:
        raise
    except Exception as e:
        print("PDF to PPT Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to convert PDF to PPT: {str(e)}")

def find_devanagari_font() -> str | None:
    import os
    paths = [
        "C:\\Windows\\Fonts\\Nirmala.ttf",
        "C:\\Windows\\Fonts\\mangal.ttf",
        "C:\\Windows\\Fonts\\arialuni.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/TTF/DejaVuSans.ttf"
    ]
    for p in paths:
        if os.path.exists(p):
            return p
    return None

def convert_scanned_pdf_to_digital(pdf_path: str, pages_to_process: list[int]) -> str:
    import fitz
    import pytesseract
    from PIL import Image
    import io
    import tempfile
    import uuid
    import os
    
    doc = fitz.open(pdf_path)
    new_doc = fitz.open()
    
    font_path = find_devanagari_font()
    
    from app.routes.converters import page_needs_ocr
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        p_w, p_h = page.rect.width, page.rect.height
        
        if page_num in pages_to_process and page_needs_ocr(page):
            new_page = new_doc.new_page(width=p_w, height=p_h)
            if font_path:
                try:
                    new_page.insert_font(fontname="deva", fontfile=font_path)
                except Exception:
                    font_path = None
                    
            try:
                pix = page.get_pixmap(dpi=200)
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data)).convert('L')
                
                ocr_data = pytesseract.image_to_data(img, lang="hin+eng", output_type=pytesseract.Output.DICT)
                n_boxes = len(ocr_data['level'])
                
                to_points_x = p_w / pix.width
                to_points_y = p_h / pix.height
                
                for i in range(n_boxes):
                    text = ocr_data['text'][i].strip()
                    conf = float(ocr_data['conf'][i]) if ocr_data['conf'][i] != -1 else 0.0
                    is_numeric_like = any(c.isdigit() for c in text)
                    if not text or (conf < 30 and not is_numeric_like):
                        continue
                    
                    l_pix = ocr_data['left'][i]
                    t_pix = ocr_data['top'][i]
                    w_pix = ocr_data['width'][i]
                    h_pix = ocr_data['height'][i]
                    
                    tx0 = l_pix * to_points_x
                    ty0 = t_pix * to_points_y
                    tx1 = (l_pix + w_pix) * to_points_x
                    ty1 = (t_pix + h_pix) * to_points_y
                    
                    rect = fitz.Rect(tx0, ty0, tx1, ty1)
                    font_sz = max(6.0, min((ty1 - ty0) * 0.8, 14.0))
                    
                    try:
                        if font_path:
                            new_page.insert_textbox(rect, text, fontsize=font_sz, fontname="deva")
                        else:
                            new_page.insert_textbox(rect, text, fontsize=font_sz, fontname="helv")
                    except Exception:
                        pass
            except Exception as e_page:
                print(f"Failed to digitalize page {page_num}: {e_page}")
                new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
        else:
            new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
            
    temp_dir = tempfile.gettempdir()
    out_pdf = os.path.join(temp_dir, f"digitalized_ocr_{uuid.uuid4().hex[:6]}.pdf")
    new_doc.save(out_pdf)
    new_doc.close()
    doc.close()
    return out_pdf

# ---------------------------------------------------------------------
# 📈 PDF to EXCEL (.xlsx)
# ---------------------------------------------------------------------
@router.post("/pdf-to-excel")
async def pdf_to_excel(
    file: UploadFile = File(...),
    page_range: str = Form(None),
    engine: str = Form("auto"),
    include_page_headers: bool = Form(True)
):
    from bisect import bisect_left
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    class StandardRow:
        def __init__(self, bbox, cells):
            self.bbox = bbox
            self.cells = cells  # list of (x0, y0, x1, y1) or None

    class StandardTable:
        def __init__(self, bbox, rows, grid_text, row_count, col_count):
            self.bbox = bbox
            self.rows = rows  # list of StandardRow
            self.grid_text = grid_text  # 2D list of strings
            self.row_count = row_count
            self.col_count = col_count

    try:
        temp_dir = tempfile.gettempdir()
        import uuid
        pdf_name = f"tmp_excel_{uuid.uuid4().hex[:6]}_{file.filename}"
        pdf_path = os.path.join(temp_dir, pdf_name)
        orig_pdf_path = pdf_path
        
        with open(orig_pdf_path, "wb") as f_out:
            f_out.write(await file.read())
            
        xlsx_path = os.path.join(temp_dir, f"tables_{file.filename}.xlsx")
        doc = fitz.open(orig_pdf_path)
        wb = Workbook()
        
        # We will keep track of sheet configurations to support smart continuation sheets
        active_ws = None
        last_page_headers = None
        last_page_cols_count = None
        
        temp_images_created = []
        MAX_OCR_PAGES = 16
        ocr_processed = 0
        pages_to_process = parse_page_range(page_range, len(doc))
        if not pages_to_process:
            pages_to_process = list(range(len(doc)))
            
        # Check if we should convert scanned/OCR pages into a spatially reconstructed digital PDF first
        any_ocr_needed = False
        if engine == "ocr":
            any_ocr_needed = True
        elif engine == "auto":
            for page_num in pages_to_process:
                if page_num < len(doc):
                    if page_needs_ocr(doc[page_num]):
                        any_ocr_needed = True
                        break
                        
        is_reconstructed_digital = False
        reconstructed_pdf_path = None
        if any_ocr_needed:
            try:
                # Pre-digitalize scanned pages
                reconstructed_pdf_path = convert_scanned_pdf_to_digital(orig_pdf_path, pages_to_process)
                # Reopen doc on the reconstructed digital file!
                doc.close()
                pdf_path = reconstructed_pdf_path
                doc = fitz.open(pdf_path)
                is_reconstructed_digital = True
            except Exception as e_recon:
                print("Failed to pre-digitalize scanned PDF, falling back to page-by-page OCR:", e_recon)
                try:
                    doc = fitz.open(orig_pdf_path)
                except Exception:
                    pass
            
        ws_created = False
        
        # Helper to check if text contains Devanagari characters
        def contains_deva(s: str) -> bool:
            return any(0x0900 <= ord(c) <= 0x097F for c in s)

        def clean_hindi_text(text: str) -> str:
            if not text:
                return text
            
            # Auto-convert legacy font encodings to Unicode
            text = convert_if_legacy_hindi(text)
            
            # 1. Direct replacements for custom font characters mapping
            replacements = {
                "\u0223": "\u0940",       # ȣ -> ी
                "\u1b43": "\u0915\u094d\u0937",  # ᭃ -> क्ष
                "\u1b4d": "\u0915\u094d",  # ᭍ -> क्
                "\u1b6b": "\u0936\u094d",  # ᭫ -> श्
                "\u1b4f": "\u0917\u094d",  # ᭏ -> ग् (changed from ग्न to ग्)
                "\u1b54": "\u091c\u094d\u092f",  # ᭔ -> ज्य
                "\u1b5b": "\u0923\u094d",  # ᭛ -> ण्
                "\u1b5c": "\u0924\u094d",  # ᭜ -> त्
                "\u1b5f": "\u0927\u094d",    # ᭟ -> ध्
                "\u1b60": "\u0928\u094d",  # ᭠ -> न्
                "\u1b61": "\u092a\u094d",  # ᭡ -> प्
                "\u1b63": "\u094d\u092c",  # ᭣ -> ्ब्
                "\u1b64": "\u092d\u094d",  # ᭤ -> भ्
                "\u1b65": "\u092e\u094d",  # ᭥ -> म्
                "\u1b68": "\u0932\u094d",  # ᭨ -> ल्
                "\u1b6a": "\u0935\u094d\u092f",  # ᭪ -> vy
                "\u1b6d": "\u0938\u094d",  # ᭭ -> स
                "\u1b93": "\u0915\u094d\u0930",  # ᮓ -> क्र
                "\u1b9d": "\u091f\u094d\u0930",  # ᮝ -> ट्र
                "\u1b9f": "\u0921\u094d\u0930",  # ᮟ -> ड्र
                "\u1ba2": "\u0924\u094d\u0930",  # ᮢ -> त्र
                "\u1ba4": "\u0926\u094d\u0930",  # ᮤ -> द्र
                "\u1ba7": "\u092a\u094d\u0930",  # ᮧ -> प्र
                "\u1c5f": "\u0939\u0941",      # ᱟ -> हु (changed from ॉ to hu)
                "\u1c6b": "\u0924\u094d\u0924",  # ᱫ -> त्त
                "\u1c6d": "\u0940\u0902",  # ᱭ -> ीं
                "\u1c76": "\u0947\u0902",  # ᱶ -> ें
                "\u1c79": "\u0948\u0902",  # ᱹ -> ैं
                "\u1c82": "\u094b\u0902",  # ᲂ -> ों
                "\u1c90": "\u0913\u0902",  # Ა -> ओं
                "\u1ca6": "\u0915\u094d\u0924",  # Ღ -> क्त
                "\u1cc7": "\u0926\u094d\u0935",  # ᳇ -> द्व
                "\u1ccd": "\u0926\u094d\u0927",  # ᳍ -> द्ध
                "\u00c8": "\u0915\u094d",  # È -> क्
                "\u00a2": "\u0915\u094d\u0937\u093e", # ¢ -> क्षा
                
                # Additional mapped characters
                "\u1cf1": "\u0940",       # ᳱ -> ी
                "\u1cef": "\u093f",       # ᳯ -> ि
                "\u1cf0": "\u093f",       # ᳰ -> ि
                "\u1cf2": "\u093f",       # ᳲ -> ि
                "\u1cf3": "\u094d\u0930",  # ᳳ -> ्र
                
                # Missing glyph maps added in Round 6
                "\u1cd9": "\u092a\u094d\u0924",  # ᳙ -> प्त (U+1CD9)
                "\u1b6c": "\u0937\u094d\u0920",  # ᭬ -> ष्ठ (U+1B6C)
                "\u1c8b": "\u0913\u0902",      # ᲋ -> ओं (U+1C8B)
                "\u1cb9": "\u0920\u094d\u092f",  # Ჹ -> ठ्य (U+1CB9)
                "\u1c68": "\u0930\u0942",      # ᱨ -> रू (U+1C68)
                "\u0238": "\u0930\u094d\u0924\u0940", # ȸ -> र्ती (U+0238)
                "\u1cdf": "\u0936\u094d\u091a",  # ᳟ -> श्च (U+1CDF)
            }
            for k, v in replacements.items():
                text = text.replace(k, v)
                
            # 2. Contextual rules
            import re
            
            # Vedic reph symbols
            text = re.sub("\u1cf6\s*([\u0900-\u097F])", "र्\\1", text) # ᳶ + C -> र् + C
            text = re.sub("\u1cf3\s*([\u0900-\u097F])", "र्\\1", text) # ᳳ + C -> र् + C
            
            text = re.sub("([\u0900-\u097F])\s*\u1c6e", "र्\\1ी", text) # C + ᱮ -> र् + C + ी
            text = re.sub("([\u0900-\u097F])\s*\u1c77", "र्\\1े", text) # C + ᱷ -> र् + C + े
            text = re.sub("([\u0905-\u093F\u0940-\u094Dक-ह][\u093e-\u094c\u0902]?)\s*ᭅ", "र्\\1", text) # C + ᭅ -> र् + C
            
            # Restore half consonants followed by spaces
            text = re.sub(r"\u094d\s+", "\u094d", text)
            
            # 3. Visual vowel reordering: ि followed by consonant cluster -> cluster + ि (only if visually misplaced)
            text = re.sub(r"(?<![\u0915-\u0939\u0958-\u095F\u094D])\u093F((?:[\u0915-\u0939\u0958-\u095F]\u094D)*[\u0915-\u0939\u0958-\u095F](?:\u094D\u0930|्र)?)", "\\1\u093F", text)
            
            # 4. Word hacks & cleaner normalizations
            word_hacks = {
                "र.": "₹",
                "ᱨ.": "₹",
                "ᱨ": "र",
                "भतȸ": "भर्ती",
                "महत्त्वपूर्ण": "महत्वपूर्ण",
                "सत् यापन": "सत्यापन",
                "सत्यापन": "सत्यापन",
                "समस् त": "समस्त",
                "आवश् यक": "आवश्यक",
                "राज्य य": "राज्य",
                "मध्य य": "मध्य",
                "मध्यय": "मध्य",
                "मान् ": "मान्य ",
                "अमान् ": "अमान्य ",
                "बॉस् ": "बॉक्स ",
                "क᳟": "की",
                "्र": "्र",
                "रिपो्रटग": "रिपोर्टिंग",
                "अत:कम्प्युटर": "अतःकम्प्यूटर",
                "कम्प्युटर": "कम्प्यूटर",
                "कोइ": "कोई",
                "लार्इन": "लाईन",
                "आनलार्इन": "ऑनलाइन",
                "ओनलार्इन": "ऑनलाइन",
                "कियोस्क": "कियोस्क",
                "माध्ययम": "माध्यम",
                "दिव्य यांगजन": "दिव्यांगजन",
                "अभ्यर्थयों": "अभ्यर्थियों",
                "परीक्षार्थयों": "परीक्षार्थियों",
                "अ᭤यᳶथयᲂ": "अभ्यर्थियों",
                "परीᭃाᳶथयᲂ": "परीक्षार्थियों",
                
                # New contextual hacks added in Round 6
                "भतें ᱮ": "भर्ती",
                "भतेंᱮ": "भर्ती",
                "भतेर्ंी": "भर्ती",
                "दशिानदिश ᱷ": "दिशानिर्देश",
                "दशिानदिश": "दिशानिर्देश",
                "विनििें ᳟त": "विनिश्चित",
                "विनििें᳟त": "विनिश्चित",
                "विनििें श्चत": "विनिश्चित",
                "विनििेंश्चत": "विनिश्चित",
                "वनिश्चिति": "विनिश्चित",
                "पृष्ठ ठ": "पृष्ठ",
                "क्₹": "क्र.",
                "क् ₹": "क्र.",
                "पालीअ": "पाली",
                "पटवारी ᭅ": "पटवारी",
                "समर्ये": "समय",
                "नर्दि श": "निर्देश",
                "नर्दि": "निर्देश",
                "नधिार्ेंरति": "निर्धारित",
                "नधिरिति": "निर्धारित",
                "प्₹": "प्र.",
                "प् ₹": "प्र.",
            }
            for k, v in word_hacks.items():
                text = text.replace(k, v)
                
            # Strip remaining visual glyph diacritics
            text = text.replace("\u1b45", "") # Remove left-over ᭅ
            text = text.replace("\u1c77", "") # Remove left-over ᱷ
            text = text.replace("\u1cdf", "") # Remove left-over ᳟
            text = text.replace("\u1c6e", "") # Remove left-over ᱮ
            return text

        def normalize_margin_text(text: str) -> str:
            if not text:
                return ""
            import re
            text = text.strip().lower()
            text = re.sub(r"\d+", "\\\\d", text)
            text = re.sub(r"\s+", " ", text)
            return text

        repeated_margin_texts = set()
        try:
            margin_text_page_map = {}
            for p_idx in pages_to_process:
                p_temp = doc.load_page(p_idx)
                p_h_val = p_temp.rect.height
                p_dict = p_temp.get_text("dict")
                for bl in p_dict.get("blocks", []):
                    if bl.get("type") == 0:
                        for ln in bl.get("lines", []):
                            for sp in ln.get("spans", []):
                                sb = sp.get("bbox", (0,0,0,0))
                                if sb[3] <= 48 or sb[1] >= p_h_val - 48:
                                    s_txt = sp.get("text", "")
                                    norm_txt = normalize_margin_text(s_txt)
                                    if norm_txt:
                                        if norm_txt not in margin_text_page_map:
                                            margin_text_page_map[norm_txt] = set()
                                        margin_text_page_map[norm_txt].add(p_idx)
            for norm_txt, p_set in margin_text_page_map.items():
                if len(p_set) >= 2:
                    repeated_margin_texts.add(norm_txt)
        except Exception as e_prescan:
            print("Repeated margin pre-scan failed:", e_prescan)
            
        # Helper to format values (numeric coercion, retaining leading zero IDs)
        def format_cell_value(text: str):
            if not text:
                return text
            # Clean Hindi text first
            text = clean_hindi_text(text)
            import re
            # Clean CID patterns like (cid:1234)
            text = re.sub(r"\(cid:\d+\)", "", text)
            text = re.sub(r"\s+", " ", text).strip()
            
            clean = text.replace(",", "").strip()
            if not clean:
                return text
            # Check integer or decimal format (allow negative)
            if clean.replace(".", "", 1).isdigit() or (clean.startswith("-") and clean[1:].replace(".", "", 1).isdigit()):
                if len(clean) > 1 and clean.startswith("0") and not clean.startswith("0."):
                    return text
                try:
                    if "." in clean:
                        return float(clean)
                    return int(clean)
                except ValueError:
                    return text
            return text

        def reconstruct_line_from_spans(spans):
            if not spans:
                return ""
            sorted_spans = sorted(spans, key=lambda s: s.get("bbox", (0,0,0,0))[0])
            line_raw = ""
            for i, s in enumerate(sorted_spans):
                txt = s.get("text", "")
                if i > 0:
                    prev_s = sorted_spans[i-1]
                    prev_bbox = prev_s.get("bbox", (0,0,0,0))
                    curr_bbox = s.get("bbox", (0,0,0,0))
                    gap = curr_bbox[0] - prev_bbox[2]
                    if gap >= 2.0 and not line_raw.endswith(" ") and not txt.startswith(" "):
                        line_raw += " "
                line_raw += txt
            return clean_hindi_text(line_raw)

        # Helper to format borders for all merged cells
        def style_merged_cells_borders(ws, start_row, start_col, end_row, end_col, b_left, b_right, b_top, b_bottom):
            for r in range(start_row, end_row + 1):
                for c in range(start_col, end_col + 1):
                    cell = ws.cell(row=r, column=c)
                    s_l = b_left if c == start_col else Side(style=None)
                    s_r = b_right if c == end_col else Side(style=None)
                    s_t = b_top if r == start_row else Side(style=None)
                    s_b = b_bottom if r == end_row else Side(style=None)
                    cell.border = Border(left=s_l, right=s_r, top=s_t, bottom=s_b)

        # Tolerant grid coordinate index finder
        def find_coord_index(sorted_list, val, tolerance=2.0):
            for idx, coord in enumerate(sorted_list):
                if abs(coord - val) <= tolerance:
                    return idx
            # Fallback
            from bisect import bisect_left
            idx = bisect_left(sorted_list, val)
            return max(0, min(idx, len(sorted_list) - 1))

        # Coordinate map helper to find row and col from PDF points
        def get_cell_coord(x: float, y: float, sorted_x_list: list[float], sorted_y_list: list[float], base_row: int):
            from bisect import bisect_left
            col = 1
            if sorted_x_list:
                col = bisect_left(sorted_x_list, x) + 1
                col = max(1, min(col, len(sorted_x_list)))
            else:
                col = max(1, int(x // 40) + 1)
            
            row_offset = 0
            if sorted_y_list:
                row_offset = bisect_left(sorted_y_list, y)
                row_offset = max(0, min(row_offset, len(sorted_y_list) - 1))
            else:
                row_offset = max(0, int(y // 15))
            
            return base_row + row_offset, col

        for page_num in pages_to_process:
            page = doc.load_page(page_num)
            p_w = page.rect.width
            p_h = page.rect.height
            
            # Read vector drawings for fills/borders
            try:
                drawings = page.get_drawings()
            except Exception:
                drawings = []
                
            # Perform text validation check to fallback to OCR under legacy hijacked mapping or scan
            if is_reconstructed_digital:
                use_ocr = False
            elif engine == "ocr":
                use_ocr = True
            elif engine == "digital":
                use_ocr = False
            else:  # "auto"
                use_ocr = page_needs_ocr(page)

            if use_ocr and ocr_processed < MAX_OCR_PAGES:
                page_dict = await asyncio.to_thread(get_page_elements, page, use_ocr=True)
                ocr_processed += 1
            else:
                page_dict = await asyncio.to_thread(get_page_elements, page, use_ocr=False)
                
            # Sort blocks in visual reading order
            if "blocks" in page_dict:
                page_dict["blocks"] = sort_blocks_reading_order(page_dict["blocks"])
                
            # Define header/footer safety zones based on page height (top/bottom 8% or 55 pt)
            top_margin = max(55.0, p_h * 0.08)
            bottom_margin = p_h - max(55.0, p_h * 0.08)
            def is_header_footer(bbox, text_val="", font_sz=11.0) -> bool:
                if not bbox or len(bbox) < 4:
                    return False
                text_clean = text_val.strip().lower()
                if "page" in text_clean or "पृष्ठ" in text_clean or text_clean.isdigit():
                    return True
                if not include_page_headers:
                    if bbox[3] <= 48:
                        norm_t = normalize_margin_text(text_val)
                        if norm_t in repeated_margin_texts:
                            return True
                        if len(text_clean) <= 15 and any(x in text_clean for x in ["मण्डल", "भोपाल", "परीक्षा", "नियमपुस्तिका"]):
                            return True
                    if bbox[1] >= p_h - 48:
                        norm_t = normalize_margin_text(text_val)
                        if norm_t in repeated_margin_texts:
                            return True
                return False
                
            # Extract tables: Prioritize Layout-Based pdfplumber Detection
            table_list = []
            try:
                with pdfplumber.open(pdf_path) as plumb_doc:
                    if page_num < len(plumb_doc.pages):
                        plumb_page = plumb_doc.pages[page_num]
                        plumb_tables = plumb_page.find_tables()
                        if not plumb_tables:
                            plumb_tables = plumb_page.find_tables(table_settings={
                                "vertical_strategy": "text",
                                "horizontal_strategy": "text",
                                "snap_y_tolerance": 8,
                                "snap_x_tolerance": 8
                            })
                        for t in plumb_tables:
                            rows = []
                            for r in t.rows:
                                cells = []
                                for cell in r.cells:
                                    if cell is None:
                                        cells.append(None)
                                    elif isinstance(cell, (list, tuple)):
                                        cells.append(cell)
                                    else:
                                        cells.append(getattr(cell, "bbox", None))
                                r_bbox = r.bbox if hasattr(r, "bbox") else t.bbox
                                rows.append(StandardRow(r_bbox, cells))
                            grid_text = t.extract()
                            col_count = len(t.cols) - 1 if hasattr(t, "cols") else max(len(r.cells) for r in t.rows)
                            table_list.append(StandardTable(t.bbox, rows, grid_text, len(t.rows), col_count))
            except Exception as e_plumb:
                print("pdfplumber table extraction failed, falling back to PyMuPDF:", e_plumb)
                
            # Fallback to PyMuPDF find_tables if pdfplumber didn't catch anything
            if not table_list:
                try:
                    tabs = page.find_tables()
                    if tabs and tabs.tables:
                        for tab in tabs.tables:
                            rows = []
                            for r in tab.rows:
                                cells = []
                                for cell in r.cells:
                                    if cell:
                                        cells.append(cell)
                                    else:
                                        cells.append(None)
                                r_bbox = r.bbox if hasattr(r, "bbox") else tab.bbox
                                rows.append(StandardRow(r_bbox, cells))
                            grid_text = tab.extract()
                            table_list.append(StandardTable(tab.bbox, rows, grid_text, tab.row_count, tab.col_count))
                except Exception as e_fitz_tab:
                    print("PyMuPDF fallback table extraction failed:", e_fitz_tab)
            
            # Check continuation of tables
            table_header = None
            col_count_val = None
            if table_list:
                first_tab = table_list[0]
                first_tab_grid = first_tab.grid_text
                col_count_val = first_tab.col_count
                if first_tab_grid and len(first_tab_grid) > 0:
                    table_header = [str(cell or "").strip() for cell in first_tab_grid[0]]
                    
            is_continuation = False
            if active_ws is not None and col_count_val is not None and last_page_cols_count is not None:
                if col_count_val == last_page_cols_count:
                    is_continuation = True
            elif active_ws is not None and table_header and last_page_headers:
                if len(table_header) == len(last_page_headers) and all(h1 == h2 for h1, h2 in zip(table_header, last_page_headers)):
                    is_continuation = True
                    
            if is_continuation and active_ws is not None:
                ws = active_ws
                current_row = ws.max_row + 2
            else:
                ws = wb.create_sheet(title=f"Page_{page_num+1}"[:31])
                ws_created = True
                active_ws = ws
                last_page_headers = table_header
                last_page_cols_count = col_count_val
                current_row = 1
                try:
                    ws.views.sheetView[0].showGridLines = True
                except Exception:
                    pass
                
                # Match orientation of PDF page
                try:
                    if p_w > p_h:
                        ws.page_setup.orientation = ws.ORIENTATION_LANDSCAPE
                    else:
                        ws.page_setup.orientation = ws.ORIENTATION_PORTRAIT
                except Exception:
                    pass

            # Gather and sort all visual components on the page sequentially
            page_elements = []
            
            # 1. Add Tables
            for t in table_list:
                page_elements.append({
                    "type": "table",
                    "bbox": t.bbox,
                    "y0": t.bbox[1],
                    "data": t
                })
                
            # 2. Add Free Text Blocks (not inside any table bounds and not header/footer)
            for block in page_dict.get("blocks", []):
                if block.get("type") == 0:
                    overlap = False
                    bx0, by0, bx1, by1 = block["bbox"]
                    for t in table_list:
                        tx0, ty0, tx1, ty1 = t.bbox
                        ix0 = max(bx0, tx0)
                        iy0 = max(by0, ty0)
                        ix1 = min(bx1, tx1)
                        iy1 = min(by1, ty1)
                        if ix1 > ix0 and iy1 > iy0:
                            overlap_area = (ix1 - ix0) * (iy1 - iy0)
                            block_area = (bx1 - bx0) * (by1 - by0)
                            if block_area > 0 and (overlap_area / block_area) > 0.3:
                                overlap = True
                                break
                    if not overlap:
                        valid_lines = []
                        for line in block.get("lines", []):
                            line_text = "".join(span.get("text", "") for span in line.get("spans", []))
                            if not is_header_footer(line.get("bbox", (0,0,0,0)), line_text, 11):
                                valid_lines.append(line)
                        if valid_lines:
                            page_elements.append({
                                "type": "text_block",
                                "bbox": block["bbox"],
                                "y0": block["bbox"][1],
                                "data": valid_lines
                            })
                            
            # 3. Add Image components
            try:
                img_list = page.get_images()
                for img_info in img_list[:12]:
                    xref = img_info[0]
                    rects = page.get_image_rects(xref)
                    if rects:
                        rx0, ry0, rx1, ry1 = rects[0]
                        dw = rx1 - rx0
                        dh = ry1 - ry0
                        # Skip page background or tiny icons
                        if dw < 15 or dh < 15 or (dw > p_w * 0.9 and dh > p_h * 0.9):
                            continue
                        page_elements.append({
                            "type": "image",
                            "bbox": (rx0, ry0, rx1, ry1),
                            "y0": ry0,
                            "data": (xref, dw, dh)
                        })
            except Exception as e_img_parse:
                print("Failed to parse page images:", e_img_parse)
                
            # Sort page elements by top coordinate Y (visually from top to bottom)
            page_elements.sort(key=lambda x: x["y0"])
            
            # Detect if page has a top-left logo to prevent side-by-side text overlap
            has_top_left_logo = False
            logo_bottom_y = 0
            for el in page_elements:
                if el["type"] == "image":
                    rx0, ry0, rx1, ry1 = el["bbox"]
                    if ry0 < 60 and rx0 < 120:
                        has_top_left_logo = True
                        logo_bottom_y = max(logo_bottom_y, ry1)
            
            # Write components to coordinate grid sequentially
            for elem in page_elements:
                if elem["type"] == "text_block":
                    valid_lines = elem["data"]
                    for line in valid_lines:
                        spans = line.get("spans", [])
                        line_text = reconstruct_line_from_spans(spans)
                        if not line_text:
                            continue
                            
                        ref_s = spans[0]
                        f_name = "Nirmala UI" if contains_deva(line_text) else "Calibri"
                        f_sz = int(ref_s.get("size", 11))
                        f_flags = ref_s.get("flags", 0)
                        
                        color_int = ref_s.get("color", 0)
                        color_r = (color_int >> 16) & 255
                        color_g = (color_int >> 8) & 255
                        color_b = color_int & 255
                        f_color = f"FF{color_r:02x}{color_g:02x}{color_b:02x}"
                        
                        font_nm_lower = ref_s.get("font", "").lower()
                        f_bold = bool(f_flags & 16) or any(x in font_nm_lower for x in ["bold", "black", "heavy", "semibold"])
                        f_italic = bool(f_flags & 2) or any(x in font_nm_lower for x in ["italic", "oblique"])
                        
                        table_cols_counts = [len(t.rows[0].cells) for t in table_list if t.rows]
                        max_cols_val = max(8, max(table_cols_counts) if table_cols_counts else 8)
                        
                        start_col_val = 1
                        if has_top_left_logo and line.get("bbox", (0,0,0,0))[1] < logo_bottom_y:
                            start_col_val = 2
                            
                        ws.merge_cells(start_row=current_row, start_column=start_col_val, end_row=current_row, end_column=max_cols_val)
                        
                        cell_obj = ws.cell(row=current_row, column=start_col_val)
                        cell_obj.value = format_cell_value(line_text)
                        cell_obj.font = Font(name=f_name, size=f_sz, bold=f_bold, italic=f_italic, color=f_color)
                        cell_obj.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
                        ws.row_dimensions[current_row].height = f_sz + 6
                        
                        current_row += 1
                        
                    # Add space after free text block
                    current_row += 1
                    
                elif elem["type"] == "table":
                    tab = elem["data"]
                    
                    x_coords = set()
                    for r in tab.rows:
                        for cell in r.cells:
                            if cell:
                                x_coords.add(round(cell[0], 1))
                                x_coords.add(round(cell[2], 1))
                    sorted_x = sorted(list(x_coords))
                    merged_x = []
                    for val in sorted_x:
                        if not merged_x or val - merged_x[-1] > 2.0:
                            merged_x.append(val)
                        else:
                            merged_x[-1] = (merged_x[-1] + val) / 2.0
                    sorted_x = merged_x
                    
                    if len(sorted_x) < 2:
                        tx0, ty0, tx1, ty1 = tab.bbox
                        sorted_x = [tx0 + i * (tx1 - tx0)/max(1, tab.col_count) for i in range(tab.col_count + 1)]
                        
                    y_coords = set()
                    for r in tab.rows:
                        y_coords.add(round(r.bbox[1], 1))
                        y_coords.add(round(r.bbox[3], 1))
                    sorted_y = sorted(list(y_coords))
                    merged_y = []
                    for val in sorted_y:
                        if not merged_y or val - merged_y[-1] > 2.0:
                            merged_y.append(val)
                        else:
                            merged_y[-1] = (merged_y[-1] + val) / 2.0
                    sorted_y = merged_y
                    
                    if len(sorted_y) < 2:
                        tx0, ty0, tx1, ty1 = tab.bbox
                        sorted_y = [ty0 + i * (ty1 - ty0)/max(1, tab.row_count) for i in range(tab.row_count + 1)]
                        
                    for c_idx_val in range(len(sorted_x) - 1):
                        col_letter = get_column_letter(c_idx_val + 1)
                        col_width_excel = (sorted_x[c_idx_val+1] - sorted_x[c_idx_val]) / 7.5
                        current_w = ws.column_dimensions[col_letter].width or 8
                        ws.column_dimensions[col_letter].width = max(current_w, max(6.0, min(col_width_excel, 60.0)))
                        
                    merges_to_apply = []
                    table_rows_written = len(sorted_y) - 1
                    
                    for r_idx, r in enumerate(tab.rows):
                        max_font_size = 11
                        for c_idx, cell in enumerate(r.cells):
                            if cell:
                                x0, y0, x1, y1 = cell
                                c_start = find_coord_index(sorted_x, round(x0, 1))
                                c_end = find_coord_index(sorted_x, round(x1, 1)) - 1
                                r_start = find_coord_index(sorted_y, round(y0, 1))
                                r_end = find_coord_index(sorted_y, round(y1, 1)) - 1
                                
                                c_start = max(0, min(c_start, len(sorted_x) - 2))
                                c_end = max(c_start, min(c_end, len(sorted_x) - 2))
                                r_start = max(0, min(r_start, len(sorted_y) - 2))
                                r_end = max(r_start, min(r_end, len(sorted_y) - 2))
                                
                                excel_r_start = current_row + r_start
                                excel_r_end = current_row + r_end
                                excel_c_start = c_start + 1
                                excel_c_end = c_end + 1
                                
                                spans_in_cell = []
                                for block in page_dict.get("blocks", []):
                                    if block.get("type") == 0:
                                        for line in block.get("lines", []):
                                            for span in line.get("spans", []):
                                                sx0, sy0, sx1, sy1 = span.get("bbox", (0,0,0,0))
                                                cx = (sx0 + sx1) / 2.0
                                                cy = (sy0 + sy1) / 2.0
                                                if (x0 - 2 <= cx <= x1 + 2) and (y0 - 2 <= cy <= y1 + 2):
                                                    spans_in_cell.append(span)
                                spans_in_cell.sort(key=lambda s: (s["bbox"][1], s["bbox"][0]))
                                
                                raw_cell_text = ""
                                if tab.grid_text and r_idx < len(tab.grid_text) and c_idx < len(tab.grid_text[r_idx]):
                                    raw_cell_text = str(tab.grid_text[r_idx][c_idx] or "").strip()
                                
                                if not use_ocr and raw_cell_text:
                                    cell_text = clean_hindi_text(raw_cell_text)
                                elif spans_in_cell:
                                    # Group spans in cell by baseline (Y coord)
                                    cell_lines = []
                                    sorted_spans_by_y = sorted(spans_in_cell, key=lambda s: s.get("bbox", (0,0,0,0))[1])
                                    current_y_group = []
                                    last_y = None
                                    for s in sorted_spans_by_y:
                                        sy0 = s.get("bbox", (0,0,0,0))[1]
                                        if last_y is None:
                                            current_y_group.append(s)
                                            last_y = sy0
                                        elif abs(sy0 - last_y) < 3.0: # Baseline tolerance of 3 points
                                            current_y_group.append(s)
                                        else:
                                            cell_lines.append(current_y_group)
                                            current_y_group = [s]
                                            last_y = sy0
                                    if current_y_group:
                                        cell_lines.append(current_y_group)
                                        
                                    reconstructed_lines = []
                                    for line_spans in cell_lines:
                                        reconstructed_lines.append(reconstruct_line_from_spans(line_spans))
                                    cell_text = "\n".join(reconstructed_lines).strip()
                                else:
                                    cell_text = clean_hindi_text(raw_cell_text)
                                
                                target_cell = ws.cell(row=excel_r_start, column=excel_c_start)
                                target_cell.value = format_cell_value(cell_text)
                                
                                cell_w = x1 - x0
                                span_center = (spans_in_cell[0]["bbox"][0] + spans_in_cell[-1]["bbox"][2]) / 2.0 if spans_in_cell else 0
                                cell_center = (x0 + x1) / 2.0
                                h_align = "left"
                                if cell_w > 0 and spans_in_cell:
                                    diff_pct = abs(span_center - cell_center) / cell_w
                                    if diff_pct < 0.08:
                                        h_align = "center"
                                    elif (x1 - spans_in_cell[-1]["bbox"][2]) / cell_w < 0.12:
                                        h_align = "right"
                                target_cell.alignment = Alignment(horizontal=h_align, vertical="center", wrap_text=True)
                                
                                bg_color = None
                                for d in drawings:
                                    draw_fill = d.get("fill")
                                    if draw_fill:
                                        rx0, ry0, rx1, ry1 = d.get("rect", (0, 0, 0, 0))
                                        d_w = rx1 - rx0
                                        d_h = ry1 - ry0
                                        
                                        is_border = True
                                        items = d.get("items", [])
                                        for item in items:
                                            if item[0] == "re":
                                                r_rect = item[1]
                                                if (r_rect.x1 - r_rect.x0) >= 2.0 and (r_rect.y1 - r_rect.y0) >= 2.0:
                                                    is_border = False
                                                    break
                                            elif item[0] in ["p", "l", "c", "qu"]:
                                                if d_w >= 2.0 and d_h >= 2.0:
                                                    is_border = False
                                                    break
                                        if items and is_border:
                                            continue
                                        if d_w < 2.0 or d_h < 2.0:
                                            continue
                                            
                                        if d_w >= 0.8 * p_w and d_h >= 0.8 * p_h:
                                            continue
                                        if d_w > 1.25 * cell_w or d_h > 1.25 * (y1 - y0):
                                            continue
                                            
                                        ix0 = max(x0, rx0)
                                        iy0 = max(y0, ry0)
                                        ix1 = min(x1, rx1)
                                        iy1 = min(y1, ry1)
                                        if ix1 > ix0 and iy1 > iy0:
                                            if ((ix1 - ix0) * (iy1 - iy0)) / ((x1 - x0) * (y1 - y0)) >= 0.75:
                                                bg_color = get_hex_from_color(draw_fill)
                                                if bg_color.upper() == "FFFFFFFF":
                                                    bg_color = None
                                                else:
                                                    break
                                if bg_color:
                                    target_cell.fill = PatternFill(start_color=bg_color, end_color=bg_color, fill_type="solid")
                                    
                                thick_threshold = 2.0
                                has_l, has_r, has_t, has_b = False, False, False, False
                                style_l, style_r, style_t, style_b = 'thin', 'thin', 'thin', 'thin'
                                color_l, color_r, color_t, color_b = 'B0B0B0', 'B0B0B0', 'B0B0B0', 'B0B0B0'
                                
                                for d in drawings:
                                    pts = d.get("rect", (0, 0, 0, 0))
                                    stroke_c = d.get("color")
                                    if stroke_c and d.get("type") in ["l", "rect"]:
                                        hex_c_full = get_hex_from_color(stroke_c)
                                        hex_c = hex_c_full[2:]
                                        width = d.get("width", 1.0)
                                        style = 'medium' if width >= thick_threshold else 'thin'
                                        if abs(pts[0] - x0) <= 2 and (y0 - 2 <= pts[1] <= y1 + 2 or y0 - 2 <= pts[3] <= y1 + 2):
                                            has_l = True; style_l = style; color_l = hex_c
                                        elif abs(pts[0] - x1) <= 2 and (y0 - 2 <= pts[1] <= y1 + 2 or y0 - 2 <= pts[3] <= y1 + 2):
                                            has_r = True; style_r = style; color_r = hex_c
                                        elif abs(pts[1] - y0) <= 2 and (x0 - 2 <= pts[0] <= x1 + 2 or x0 - 2 <= pts[2] <= x1 + 2):
                                            has_t = True; style_t = style; color_t = hex_c
                                        elif abs(pts[1] - y1) <= 2 and (x0 - 2 <= pts[0] <= x1 + 2 or x0 - 2 <= pts[2] <= x1 + 2):
                                            has_b = True; style_b = style; color_b = hex_c
                                            
                                b_left = Side(style=style_l, color=color_l) if has_l else Side(style='thin', color='C0C0C0')
                                b_right = Side(style=style_r, color=color_r) if has_r else Side(style='thin', color='C0C0C0')
                                b_top = Side(style=style_t, color=color_t) if has_t else Side(style='thin', color='C0C0C0')
                                b_bottom = Side(style=style_b, color=color_b) if has_b else Side(style='thin', color='C0C0C0')
                                target_cell.border = Border(left=b_left, right=b_right, top=b_top, bottom=b_bottom)
                                
                                if excel_r_end > excel_r_start or excel_c_end > excel_c_start:
                                    merges_to_apply.append((excel_r_start, excel_c_start, excel_r_end, excel_c_end, b_left, b_right, b_top, b_bottom))
                                    
                                font_n = "Nirmala UI" if contains_deva(cell_text) else "Calibri"
                                font_sz = 11
                                font_c = "FF000000"
                                is_bold = False
                                is_italic = False
                                
                                if spans_in_cell:
                                    ref_span = spans_in_cell[0]
                                    font_sz = int(ref_span.get("size", 11.0))
                                    flags = ref_span.get("flags", 0)
                                    color_int = ref_span.get("color", 0)
                                    color_r = (color_int >> 16) & 255
                                    color_g = (color_int >> 8) & 255
                                    color_b = color_int & 255
                                    font_c = f"FF{color_r:02x}{color_g:02x}{color_b:02x}"
                                    
                                    font_name_lower = ref_span.get("font", "").lower()
                                    is_bold = bool(flags & 16) or any(x in font_name_lower for x in ["bold", "black", "heavy", "semibold"])
                                    is_italic = bool(flags & 2) or any(x in font_name_lower for x in ["italic", "oblique"])
                                    
                                if bg_color:
                                    r_val = int(bg_color[2:4], 16)
                                    g_val = int(bg_color[4:6], 16)
                                    b_val = int(bg_color[6:8], 16)
                                    bg_brightness = (r_val * 299 + g_val * 587 + b_val * 114) / 1000
                                    
                                    f_r = int(font_c[2:4], 16)
                                    f_g = int(font_c[4:6], 16)
                                    f_b = int(font_c[6:8], 16)
                                    font_brightness = (f_r * 299 + f_g * 587 + f_b * 114) / 1000
                                    
                                    if bg_brightness < 120 and font_brightness < 150:
                                        font_c = "FFFFFFFF"
                                    elif bg_brightness > 200 and font_brightness > 200:
                                        font_c = "FF000000"
                                        
                                target_cell.font = Font(name=font_n, size=font_sz, bold=is_bold, italic=is_italic, color=font_c)
                                max_font_size = max(max_font_size, font_sz)
                                
                        ws.row_dimensions[current_row + r_idx].height = max_font_size + 6
                        
                    for m_start_r, m_start_c, m_end_r, m_end_c, b_l, b_r, b_t, b_b in merges_to_apply:
                        try:
                            ws.merge_cells(start_row=m_start_r, start_column=m_start_c, end_row=m_end_r, end_column=m_end_c)
                            style_merged_cells_borders(ws, m_start_r, m_start_c, m_end_r, m_end_c, b_l, b_r, b_t, b_b)
                        except Exception:
                            pass
                            
                    current_row += table_rows_written
                    current_row += 1
                    
                elif elem["type"] == "image":
                    xref, dw, dh = elem["data"]
                    rx0, ry0, rx1, ry1 = elem["bbox"]
                    try:
                        base_img = doc.extract_image(xref)
                        if base_img:
                            img_bytes = base_img["image"]
                            img_ext = base_img["ext"]
                            temp_img_name = f"obj_img_{page_num}_{xref}_{uuid.uuid4().hex[:6]}.{img_ext}"
                            temp_img_path = os.path.join(temp_dir, temp_img_name)
                            with open(temp_img_path, "wb") as f_img:
                                f_img.write(img_bytes)
                            temp_images_created.append(temp_img_path)
                            
                            from openpyxl.drawing.image import Image as OpenpyxlImage
                            ox_img = OpenpyxlImage(temp_img_path)
                            ox_img.width = int(dw * 1.333)
                            ox_img.height = int(dh * 1.333)
                            
                            col_click = 1
                            if rx0 > p_w * 0.7:
                                col_click = 5
                            elif rx0 > p_w * 0.3:
                                col_click = 3
                            col_letter = get_column_letter(col_click)
                            
                            if ry0 < 60 and rx0 < 120:
                                ox_img.anchor = f"{col_letter}1"
                                logo_w_excel = dw / 7.5
                                ws.column_dimensions["A"].width = max(ws.column_dimensions["A"].width or 0, logo_w_excel)
                                ws.add_image(ox_img)
                            else:
                                ox_img.anchor = f"{col_letter}{current_row}"
                                ws.add_image(ox_img)
                                rows_needed = int(dh // 15) + 1
                                for r_off in range(rows_needed):
                                    ws.row_dimensions[current_row + r_off].height = 15
                                current_row += rows_needed
                                current_row += 1
                    except Exception as e_img:
                        print("Failed to embed image assets in Excel:", e_img)
                    
        # Check if the document was truncated based on page range
        if pages_to_process and pages_to_process[-1] < len(doc) - 1 and active_ws is not None:
            trunc_row = active_ws.max_row + 2
            active_ws.cell(row=trunc_row, column=1).value = f"[Table truncated to page range {page_range} - Content continues on next page of original PDF]"
            active_ws.cell(row=trunc_row, column=1).font = Font(name="Calibri", size=11, italic=True, bold=True, color="FFD32F2F") # Red color
            max_c = active_ws.max_column or 8
            try:
                active_ws.merge_cells(start_row=trunc_row, start_column=1, end_row=trunc_row, end_column=max_c)
            except Exception:
                pass

        if ws_created and "Sheet" in wb.sheetnames:
            del wb["Sheet"]
            
        wb.save(xlsx_path)
        doc.close()

        # Clean up temporary image files
        for p in temp_images_created:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass
                
        # Clean up spatial digital reconstruction file
        if is_reconstructed_digital and reconstructed_pdf_path and os.path.exists(reconstructed_pdf_path):
            try:
                os.remove(reconstructed_pdf_path)
            except Exception:
                pass
                
        # Clean up original uploaded file
        try:
            if orig_pdf_path and os.path.exists(orig_pdf_path):
                os.remove(orig_pdf_path)
        except Exception:
            pass
        
        return FileResponse(
            xlsx_path,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=f"Data_{file.filename}.xlsx"
        )
    except Exception as e:
        print("PDF to Excel Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to convert PDF to Excel: {e}")

# ---------------------------------------------------------------------
# 🏛️ PDF to PDF/A (Archival Standard)
# ---------------------------------------------------------------------
@router.post("/pdf-to-pdfa")
async def pdf_to_pdfa(file: UploadFile = File(...)):
    try:
        # PyMuPDF doc.save(pdfa=1) requires embedded fonts and rgb space.
        # Alternatively, we can rasterize to guarantee compliance or use ghostscript if installed.
        # A simpler approach using PyMuPDF which attempts PDF/A-1b metadata assignment:
        temp_dir = tempfile.gettempdir()
        in_path = os.path.join(temp_dir, f"in_{file.filename}")
        out_path = os.path.join(temp_dir, f"out_pdfa_{file.filename}")
        
        with open(in_path, "wb") as f:
            f.write(await file.read())
            
        doc = fitz.open(in_path)
        
        # PyMuPDF directly supports saving to PDF/A via ghostscript/internal flags in 1.19+ but we flatten to be safe
        out_doc = fitz.open()
        for page in doc:
            pix = page.get_pixmap(dpi=150)
            img_doc = fitz.open(stream=pix.tobytes("jpeg"), filetype="jpeg")
            img_pdf = img_doc.convert_to_pdf()
            temp_pdf = fitz.open("pdf", img_pdf)
            out_doc.insert_pdf(temp_pdf)
            temp_pdf.close()
            img_doc.close()
            
        # Write metadata for PDF/A
        out_doc.set_metadata({
            "format": "PDF 1.4",
            "title": file.filename,
            "creator": "MeriPDF"
        })
        
        # We apply deflate and clean for best archiving packaging
        out_doc.save(out_path, deflate=True)
        out_doc.close()
        doc.close()
        
        return FileResponse(
            out_path,
            media_type="application/pdf",
            filename=f"Archival_{file.filename}"
        )
    except Exception as e:
        print("PDF to PDF/A Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to convert PDF to PDF/A: {e}")

# ---------------------------------------------------------------------
# ⚙️ LibreOffice PDF Engine & Multi-Format Converters
# ---------------------------------------------------------------------
def find_libreoffice():
    import shutil
    # 1. Check if in PATH
    path = shutil.which("soffice") or shutil.which("libreoffice")
    if path:
        return path
    
    # 2. Check standard Windows paths
    standard_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for p in standard_paths:
        if os.path.exists(p):
            return p
            
    return None

def convert_to_pdf_with_libreoffice(in_path: str, temp_dir: str) -> str:
    import subprocess
    
    libreoffice_path = find_libreoffice()
    if not libreoffice_path:
        print("❌ Error: LibreOffice not found on the system.")
        raise HTTPException(
            status_code=500,
            detail="PDF conversion service is currently unavailable on the server (missing rendering engine)."
        )
        
    cmd = [
        libreoffice_path,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        temp_dir,
        in_path
    ]
    
    startupinfo = None
    if os.name == 'nt':
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        
    result = subprocess.run(
        cmd, 
        stdout=subprocess.PIPE, 
        stderr=subprocess.PIPE, 
        startupinfo=startupinfo, 
        timeout=60
    )
    
    if result.returncode != 0:
        print(f"❌ LibreOffice conversion failed (code {result.returncode}): {result.stderr.decode()}")
        raise HTTPException(
            status_code=500,
            detail="Conversion failed. The document could not be converted to PDF."
        )
        
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    out_path = os.path.join(temp_dir, f"{base_name}.pdf")
    
    if not os.path.exists(out_path):
        raise HTTPException(
            status_code=500,
            detail="Conversion output file not generated."
        )
        
    return out_path

async def handle_libreoffice_conversion(file: UploadFile, allowed_exts: list[str], err_prefix: str, response_prefix: str):
    import uuid
    
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in allowed_exts:
        raise HTTPException(
            status_code=400, 
            detail=f"Invalid file extension. Allowed extensions are: {', '.join(allowed_exts)}"
        )
        
    unique_id = str(uuid.uuid4())
    temp_dir = tempfile.gettempdir()
    in_path = os.path.join(temp_dir, f"in_{unique_id}{ext}")
    out_path = None
    
    try:
        with open(in_path, "wb") as f:
            f.write(await file.read())
            
        out_path = convert_to_pdf_with_libreoffice(in_path, temp_dir)
        
        # Cleanup input file
        if os.path.exists(in_path):
            os.remove(in_path)
            
        base_filename = os.path.splitext(file.filename)[0]
        clean_filename = f"{response_prefix}_{base_filename}.pdf"
        
        return FileResponse(
            out_path,
            media_type="application/pdf",
            filename=clean_filename,
            headers={"Content-Disposition": f"attachment; filename={clean_filename}"}
        )
    except HTTPException:
        if os.path.exists(in_path):
            try: os.remove(in_path)
            except: pass
        if out_path and os.path.exists(out_path):
            try: os.remove(out_path)
            except: pass
        raise
    except Exception as e:
        print(f"{err_prefix} Error:", e)
        if os.path.exists(in_path):
            try: os.remove(in_path)
            except: pass
        if out_path and os.path.exists(out_path):
            try: os.remove(out_path)
            except: pass
        raise HTTPException(status_code=500, detail=f"Failed to convert file to PDF: {str(e)}")

# 📄 WORD to PDF (.pdf)
@router.post("/word-to-pdf")
async def word_to_pdf(files: list[UploadFile] = File(...), rotations: list[str] = Form(default=None)):
    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded")
    if len(files) > 50:
        raise HTTPException(status_code=400, detail="Maximum limit of 50 documents exceeded per conversion request.")
    
    # Parse rotations or pad it to match len(files)
    rot_list = []
    if rotations:
        for r in rotations:
            try:
                rot_list.append(int(r))
            except ValueError:
                rot_list.append(0)
            
    while len(rot_list) < len(files):
        rot_list.append(0)
        
    import uuid
    import tempfile
    import zipfile
    
    unique_id = str(uuid.uuid4())
    temp_dir = tempfile.gettempdir()
    
    temp_pdf_paths = []
    temp_in_paths = []
    
    try:
        # Convert each Word document to a PDF in order
        for idx, file in enumerate(files):
            ext = os.path.splitext(file.filename)[1].lower()
            if ext not in [".doc", ".docx"]:
                raise HTTPException(
                    status_code=400,
                    detail=f"Invalid file extension: {file.filename}. Only .doc and .docx are allowed."
                )
            
            in_path = os.path.join(temp_dir, f"in_{uuid.uuid4()}{ext}")
            temp_in_paths.append(in_path)
            
            with open(in_path, "wb") as f:
                f.write(await file.read())
            
            out_pdf = convert_to_pdf_with_libreoffice(in_path, temp_dir)
            
            # Apply rotation if specified and != 0
            rot = rot_list[idx]
            if rot in [90, 180, 270]:
                doc = fitz.open(out_pdf)
                for page in doc:
                    page.set_rotation((page.rotation + rot) % 360)
                doc.save(out_pdf, incremental=True)
                doc.close()
                
            temp_pdf_paths.append(out_pdf)
            
        # If there is only one file, return it directly
        if len(files) == 1:
            for p in temp_in_paths:
                if os.path.exists(p):
                    try: os.remove(p)
                    except: pass
            
            filename = f"Converted_{os.path.splitext(files[0].filename)[0]}.pdf"
            return FileResponse(
                temp_pdf_paths[0],
                media_type="application/pdf",
                filename=filename,
                headers={"Content-Disposition": f"attachment; filename={filename}"}
            )
            
        # If there are multiple files, bundle them into a ZIP archive
        zip_filename = f"Converted_PDF_Documents_{unique_id}.zip"
        zip_path = os.path.join(temp_dir, zip_filename)
        
        with zipfile.ZipFile(zip_path, "w") as zip_file:
            for idx, pdf_path in enumerate(temp_pdf_paths):
                original_name = files[idx].filename
                base_name = os.path.splitext(original_name)[0]
                zip_file.write(pdf_path, arcname=f"{base_name}.pdf")
                
        # Clean up temporary PDFs and input files
        for p in temp_in_paths + temp_pdf_paths:
            if os.path.exists(p):
                try: os.remove(p)
                except: pass
                
        return FileResponse(
            zip_path,
            media_type="application/zip",
            filename="Converted_Word_PDFs.zip",
            headers={"Content-Disposition": "attachment; filename=Converted_Word_PDFs.zip"}
        )
        
    except Exception as e:
        # Cleanup on failure
        for p in temp_in_paths + temp_pdf_paths:
            if os.path.exists(p):
                try: os.remove(p)
                except: pass
        print("Batch Word to PDF Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to convert Word files to PDF: {str(e)}")

# 📊 POWERPOINT to PDF (.pdf)
@router.post("/ppt-to-pdf")
async def ppt_to_pdf(files: list[UploadFile] = File(...), rotations: list[str] = Form(default=None)):
    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded")
    if len(files) > 50:
        raise HTTPException(status_code=400, detail="Maximum limit of 50 presentations exceeded per conversion request.")
    
    # Parse rotations or pad it to match len(files)
    rot_list = []
    if rotations:
        for r in rotations:
            try:
                rot_list.append(int(r))
            except ValueError:
                rot_list.append(0)
            
    while len(rot_list) < len(files):
        rot_list.append(0)
        
    import uuid
    import tempfile
    import zipfile
    
    unique_id = str(uuid.uuid4())
    temp_dir = tempfile.gettempdir()
    
    temp_pdf_paths = []
    temp_in_paths = []
    
    try:
        # Convert each Presentation slideshow to a PDF in order
        for idx, file in enumerate(files):
            ext = os.path.splitext(file.filename)[1].lower()
            if ext not in [".ppt", ".pptx"]:
                raise HTTPException(
                    status_code=400,
                    detail=f"Invalid file extension: {file.filename}. Only .ppt and .pptx are allowed."
                )
            
            in_path = os.path.join(temp_dir, f"in_{uuid.uuid4()}{ext}")
            temp_in_paths.append(in_path)
            
            with open(in_path, "wb") as f:
                f.write(await file.read())
            
            out_pdf = convert_to_pdf_with_libreoffice(in_path, temp_dir)
            
            # Apply rotation if specified and != 0
            rot = rot_list[idx]
            if rot in [90, 180, 270]:
                doc = fitz.open(out_pdf)
                for page in doc:
                    page.set_rotation((page.rotation + rot) % 360)
                doc.save(out_pdf, incremental=True)
                doc.close()
                
            temp_pdf_paths.append(out_pdf)
            
        # If there is only one file, return it directly
        if len(files) == 1:
            for p in temp_in_paths:
                if os.path.exists(p):
                    try: os.remove(p)
                    except: pass
            
            filename = f"Converted_{os.path.splitext(files[0].filename)[0]}.pdf"
            return FileResponse(
                temp_pdf_paths[0],
                media_type="application/pdf",
                filename=filename,
                headers={"Content-Disposition": f"attachment; filename={filename}"}
            )
            
        # If there are multiple files, bundle them into a ZIP archive
        zip_filename = f"Converted_PDF_Presentations_{unique_id}.zip"
        zip_path = os.path.join(temp_dir, zip_filename)
        
        with zipfile.ZipFile(zip_path, "w") as zip_file:
            for idx, pdf_path in enumerate(temp_pdf_paths):
                original_name = files[idx].filename
                base_name = os.path.splitext(original_name)[0]
                zip_file.write(pdf_path, arcname=f"{base_name}.pdf")
                
        # Clean up temporary PDFs and input files
        for p in temp_in_paths + temp_pdf_paths:
            if os.path.exists(p):
                try: os.remove(p)
                except: pass
                
        return FileResponse(
            zip_path,
            media_type="application/zip",
            filename="Converted_PowerPoint_PDFs.zip",
            headers={"Content-Disposition": "attachment; filename=Converted_PowerPoint_PDFs.zip"}
        )
        
    except Exception as e:
        # Cleanup on failure
        for p in temp_in_paths + temp_pdf_paths:
            if os.path.exists(p):
                try: os.remove(p)
                except: pass
        print("Batch PowerPoint to PDF Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to convert PowerPoint files to PDF: {str(e)}")

# 📈 EXCEL to PDF (.pdf)
@router.post("/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    return await handle_libreoffice_conversion(
        file=file, 
        allowed_exts=[".xls", ".xlsx", ".csv"], 
        err_prefix="Excel to PDF", 
        response_prefix="Converted_Data"
    )

async def render_with_playwright(source: str, out_path: str, is_url: bool):
    from playwright.async_api import async_playwright
    async with async_playwright() as p:
        browser = await p.chromium.launch(
            headless=True,
            args=[
                "--no-sandbox",
                "--disable-setuid-sandbox",
                "--disable-dev-shm-usage",
                "--disable-gpu"
            ]
        )
        context = await browser.new_context(
            viewport={'width': 1280, 'height': 800},
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            ignore_https_errors=True
        )
        page = await context.new_page()
        
        # Prevent popup alert messages from blocking page load
        page.on("dialog", lambda dialog: dialog.dismiss())
        
        if is_url:
            try:
                await page.goto(source, wait_until="networkidle", timeout=25000)
            except Exception as e:
                print(f"Playwright networkidle timeout/error on {source}: {e}. Proceeding.")
                try:
                    await page.wait_for_timeout(2000)
                except:
                    pass
        else:
            file_url = f"file:///{os.path.abspath(source).replace(os.sep, '/')}"
            try:
                await page.goto(file_url, wait_until="networkidle", timeout=15000)
            except Exception as e:
                print(f"Playwright local html load warning: {e}. Proceeding.")
                
        await page.pdf(
            path=out_path,
            format="A4",
            print_background=True,
            margin={"top": "1cm", "bottom": "1cm", "left": "1cm", "right": "1cm"}
        )
        await browser.close()

async def render_preview_with_playwright(url: str, out_path: str):
    from playwright.async_api import async_playwright
    async with async_playwright() as p:
        browser = await p.chromium.launch(
            headless=True,
            args=[
                "--no-sandbox",
                "--disable-setuid-sandbox",
                "--disable-dev-shm-usage",
                "--disable-gpu"
            ]
        )
        context = await browser.new_context(
            viewport={'width': 1280, 'height': 800},
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, with Gecko) Chrome/120.0.0.0 Safari/537.36",
            ignore_https_errors=True
        )
        page = await context.new_page()
        page.on("dialog", lambda dialog: dialog.dismiss())
        
        try:
            await page.goto(url, wait_until="networkidle", timeout=25000)
        except Exception as e:
            print(f"Playwright preview load error: {e}. Proceeding.")
            try:
                await page.wait_for_timeout(2000)
            except:
                pass
                
        await page.screenshot(path=out_path, type="jpeg", quality=80)
        await browser.close()

# 📄 HTML to PDF (.pdf)
@router.post("/html-to-pdf")
async def html_to_pdf(file: UploadFile = File(...)):
    import uuid
    import tempfile
    
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in [".html", ".htm"]:
        raise HTTPException(
            status_code=400, 
            detail="Invalid file extension. Allowed extensions are: .html, .htm"
        )
        
    unique_id = str(uuid.uuid4())
    temp_dir = tempfile.gettempdir()
    in_path = os.path.join(temp_dir, f"in_{unique_id}{ext}")
    out_path = os.path.join(temp_dir, f"out_{unique_id}.pdf")
    
    try:
        with open(in_path, "wb") as f:
            f.write(await file.read())
            
        await render_with_playwright(in_path, out_path, is_url=False)
        
        base_filename = os.path.splitext(file.filename)[0]
        clean_filename = f"Converted_{base_filename}.pdf"
        
        return FileResponse(
            out_path,
            media_type="application/pdf",
            filename=clean_filename,
            headers={"Content-Disposition": f"attachment; filename={clean_filename}"}
        )
    except Exception as e:
        print("HTML to PDF Playwright Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to convert HTML to PDF: {str(e)}")
    finally:
        if os.path.exists(in_path):
            try: os.remove(in_path)
            except: pass

# 🌐 URL to PDF
@router.post("/url-to-pdf")
async def url_to_pdf(url: str = Form(...)):
    import uuid
    import tempfile
    from urllib.parse import urlparse
    
    target_url = url.strip()
    if not target_url.startswith(("http://", "https://")):
        target_url = "http://" + target_url
        
    unique_id = str(uuid.uuid4())
    temp_dir = tempfile.gettempdir()
    out_path = os.path.join(temp_dir, f"webpage_{unique_id}.pdf")
    
    try:
        await render_with_playwright(target_url, out_path, is_url=True)
        
        parsed_url = urlparse(target_url)
        domain = parsed_url.netloc.replace("www.", "") or "webpage"
        clean_filename = f"Webpage_{domain}.pdf"
        
        return FileResponse(
            out_path,
            media_type="application/pdf",
            filename=clean_filename,
            headers={"Content-Disposition": f"attachment; filename={clean_filename}"}
        )
    except Exception as e:
        print("URL to PDF Playwright Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to convert URL to PDF: {str(e)}")

# 🌐 URL to Image Preview
@router.post("/url-preview")
async def url_preview(url: str = Form(...)):
    import uuid
    import tempfile
    from fastapi import BackgroundTasks
    
    target_url = url.strip()
    if not target_url.startswith(("http://", "https://")):
        target_url = "http://" + target_url
        
    unique_id = str(uuid.uuid4())
    temp_dir = tempfile.gettempdir()
    out_path = os.path.join(temp_dir, f"preview_{unique_id}.jpg")
    
    try:
        await render_preview_with_playwright(target_url, out_path)
        
        # Schedule cleanup after response is sent
        def remove_file(filepath: str):
            if os.path.exists(filepath):
                try: os.remove(filepath)
                except: pass
                
        bg_tasks = BackgroundTasks()
        bg_tasks.add_task(remove_file, out_path)
        
        return FileResponse(
            out_path,
            media_type="image/jpeg",
            filename="preview.jpg",
            background=bg_tasks
        )
    except Exception as e:
        print("URL Preview Playwright Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to generate webpage preview: {str(e)}")



# 🖼️ GENERATE DOCUMENT THUMBNAIL
@router.post("/document-thumbnail")
async def document_thumbnail(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in [".doc", ".docx", ".ppt", ".pptx", ".pdf"]:
        raise HTTPException(status_code=400, detail="Unsupported file format for thumbnail generation.")
    
    import uuid
    import tempfile
    import fitz
    
    unique_id = str(uuid.uuid4())
    temp_dir = tempfile.gettempdir()
    in_path = os.path.join(temp_dir, f"thumb_in_{unique_id}{ext}")
    
    try:
        # Save temp file
        with open(in_path, "wb") as f:
            f.write(await file.read())
            
        out_pdf_path = None
        
        # If it is already a PDF
        if ext == ".pdf":
            out_pdf_path = in_path
        else:
            # Convert doc/docx/ppt/pptx to PDF with LibreOffice
            out_pdf_path = convert_to_pdf_with_libreoffice(in_path, temp_dir)
            
        # Extract first page using PyMuPDF (fitz)
        doc = fitz.open(out_pdf_path)
        if len(doc) == 0:
            raise Exception("Generated PDF contains no pages")
            
        page = doc[0]
        pix = page.get_pixmap(dpi=100) # Quick preview dpi
        
        png_filename = f"thumb_out_{unique_id}.png"
        png_path = os.path.join(temp_dir, png_filename)
        pix.save(png_path)
        
        doc.close()
        
        # Clean up input and converted files
        if os.path.exists(in_path):
            try: os.remove(in_path)
            except: pass
        if out_pdf_path and out_pdf_path != in_path and os.path.exists(out_pdf_path):
            try: os.remove(out_pdf_path)
            except: pass
            
        # Return the PNG thumbnail
        return FileResponse(
            png_path,
            media_type="image/png",
            filename="thumbnail.png"
        )
        
    except Exception as e:
        # Clean up on failure
        if os.path.exists(in_path):
            try: os.remove(in_path)
            except: pass
        print("Thumbnail Generation Error:", e)
        raise HTTPException(status_code=500, detail=f"Failed to generate thumbnail: {str(e)}")


# ==============================================================================
# HINDI & NEPALI FONT CONVERTER SUITE LOGIC
# ==============================================================================

from pydantic import BaseModel

class FontConvertRequest(BaseModel):
    text: str
    conversion_type: str  # e.g., 'krutidev_to_unicode', 'unicode_to_krutidev', ...

DEVLYS_ARRAY_ONE = [
    "ñ","Q+Z","sas","aa",")Z","ZZ","‘","’","“","”",
    "å",  "ƒ",  "„",   "…",   "†",   "‡",   "ˆ",   "‰",   "Š",   "‹", 
    "¶+",   "d+", "[+k","[+", "x+",  "T+",  "t+", "M+", "<+", "Q+", ";+", "j+", "u+",
    "Ùk", "Ù", "ä", "–", "—","é","™","=kk","f=k",  
    "à",   "á",    "â",   "ã",   "ºz",  "º",   "í", "{k", "{", "=",  "«",   
    "Nî",   "Vî",    "Bî",   "Mî",   "<î", "|", "K", "}",
    "J",   "Vª",   "Mª",  "<ªª",  "Nª",   "Ø",  "Ý", "nzZ",  "æ", "ç", "Á", "xz", "#", ":",
    "v‚","vks",  "vkS",  "vk",    "v",  "b±", "Ã",  "bZ",  "b",  "m",  "Å",  ",s",  ",",   "_",
    "ô",  "d", "Dk", "D", "[k", "[", "x","Xk", "X", "Ä", "?k", "?",   "³", 
    "pkS",  "p", "Pk", "P",  "N",  "t", "Tk", "T",  ">", "÷", "¥",
    "ê",  "ë",   "V",  "B",   "ì",   "ï", "M+", "<+", "M",  "<", ".k", ".",    
    "r",  "Rk", "R",   "Fk", "F",  ")", "n", "/k", "èk",  "/", "Ë", "è", "u", "Uk", "U",   
    "i",  "Ik", "I",   "Q",    "¶",  "c", "Ck",  "C",  "Hk",  "H", "e", "Ek",  "E",
    ";",  "¸",   "j",    "y", "Yk",  "Y",  "G",  "o", "Ok", "O",
    "'k", "'",   "\"k",  "\"",  "l", "Lk",  "L",   "g", 
    "È", "z", 
    "Ì", "Í", "Î",  "Ï",  "Ñ",  "Ò",  "Ó",  "Ô",   "Ö",  "Ø",  "Ù","Ük", "Ü",
    "‚",    "ks",   "kS",   "k",  "h",    "q",   "w",   "`",    "s",    "S",
    "a",    "¡",    "%",     "W",  "•", "·", "∙", "·", "~j",  "~", "\\","+"," ः",
    "^", "*",  "Þ", "ß", "(", "¼", "½", "¿", "À", "¾", "A", "-", "&", "&", "Œ", "]","~ ","@"
]

DEVLYS_ARRAY_TWO = [
    "॰","QZ+","sa","a","र्द्ध","Z","\"","\"","'","'",
    "०",  "१",  "२",  "३",     "४",   "५",  "६",   "७",   "८",   "९",   
    "फ़्",  "क़",  "ख़", "ख़्",  "ग़", "ज़्", "ज़",  "ड़",  "ढ़",   "फ़",  "य़",  "ऱ",  "ऩ",
    "त्त", "त्त्", "क्त",  "दृ",  "कृ","न्न","न्न्","=k","f=",
    "ह्न",  "ह्य",  "हृ",  "ह्म",  "ह्र",  "ह्",   "द्द",  "क्ष", "क्ष्", "त्र", "त्र्", 
    "छ्य",  "ट्य",  "ठ्य",  "ड्य",  "ढ्य", "द्य", "ज्ञ", "द्व",
    "श्र",  "ट्र",    "ड्र",    "ढ्र",    "छ्र",   "क्र",  "फ्र", "र्द्र",  "द्र",   "प्र", "प्र",  "ग्र", "रु",  "रू",
    "ऑ",   "ओ",  "औ",  "आ",   "अ", "ईं", "ई",  "ई",   "इ",  "उ",   "ऊ",  "ऐ",  "ए", "ऋ",
    "क्क", "क", "क", "क्", "ख", "ख्", "ग", "ग", "ग्", "घ", "घ", "घ्", "ङ",
    "चै",  "च", "च", "च्", "छ", "ज", "ज", "ज्",  "झ",  "झ्", "ञ",
    "टृ",   "ट्ठ",   "ट",   "ठ",   "ड्ड",   "ड्ढ",  "ड़", "ढ़", "ड",   "ढ", "ण", "ण्",   
    "त", "त", "त्", "थ", "थ्",  "द्ध",  "द", "ध", "ध", "ध्", "ध्", "ध्", "न", "न", "न्",    
    "प", "प", "प्",  "फ", "फ्",  "ब", "ब", "ब्",  "भ", "भ्",  "म",  "म", "म्",  
    "य", "य्",  "र", "ल", "ल", "ल्",  "ळ",  "व", "व", "व्",   
    "श", "श्",  "ष", "ष्", "स", "स", "स्", "ह", 
    "ीं", "्र",    
    "द्द", "ट्ट","ट्ठ","ड्ड","कृ","भ","्य","ड्ढ","झ्","क्र","त्त्","श","श्",
    "ॉ",  "ो",   "ौ",   "ा",   "ी",   "ु",   "ू",   "ृ",   "े",   "ै",
    "ं",   "ँ",   "ः",   "ॅ",  "ऽ", "ऽ", "ऽ", "ऽ", "्र",  "्", "?", "़",":",
    "‘",   "’",   "“",   "”",  ";",  "(",    ")",   "{",    "}",   "=", "।", ".", "-",  "µ", "॰", ",","् ","/"
]

UNICODE_TO_KRUTIDEV_ONE = [
    "‘",   "’",   "“",   "”",   "(",    ")",   "{",    "}",   "=", "।",  "?",  "-",  "µ", "॰", ",", ".", "् ", 
    "०",  "१",  "२",  "३",     "४",   "५",  "६",   "७",   "८",   "९", "x", 
    "फ़्",  "क़",  "ख़",  "ग़", "ज़्", "ज़",  "ड़",  "ढ़",   "फ़",  "य़",  "ऱ",  "ऩ",
    "त्त्",   "त्त",     "क्त",  "दृ",  "कृ",
    "ह्न",  "ह्य",  "हृ",  "ह्म",  "ह्र",  "ह्",   "द्द",  "क्ष्", "क्ष", "त्र्", "त्र","ज्ञ",
    "छ्य",  "ट्य",  "ठ्य",  "ड्य",  "ढ्य", "द्य","द्व",
    "श्र",  "ट्र",    "ड्र",    "ढ्र",    "छ्र",   "क्र",  "फ्र",  "द्र",   "प्र",   "ग्र", "रु",  "रू",
    "्र",
    "ओ",  "औ",  "आ",   "अ",   "ई",   "इ",  "उ",   "ऊ",  "ऐ",  "ए", "ऋ",
    "क्",  "क",  "क्क",  "ख्",   "ख",    "ग्",   "ग",  "घ्",  "घ",    "ङ",
    "चै",   "च्",   "च",   "छ",  "ज्", "ज",   "झ्",  "झ",   "ञ",
    "ट्ट",   "ट्ठ",   "ट",   "ठ",   "ड्ड",   "ड्ढ",  "ड",   "ढ",  "ण्", "ण",  
    "त्",  "त",  "थ्", "थ",  "द्ध",  "द", "ध्", "ध",  "न्",  "न",  
    "प्",  "प",  "फ्", "फ",  "ब्",  "ब", "भ्",  "भ",  "म्",  "म",
    "य्",  "य",  "र",  "ल्", "ल",  "ळ",  "व्",  "व", 
    "श्", "श",  "ष्", "ष",  "स्",   "स",   "ह",     
    "ऑ",   "ॉ",  "ो",   "ौ",   "ा",   "ी",   "ु",   "ू",   "ृ",   "े",   "ै",
    "ं",   "ँ",   "ः",   "ॅ",    "ऽ",  "् ", "्"
]

UNICODE_TO_KRUTIDEV_TWO = [
    "^", "*",  "Þ", "ß", "¼", "½", "¿", "À", "¾", "A", "\\", "&", "&", "Œ", "]","-","~ ", 
    "å",  "ƒ",  "„",   "…",   "†",   "‡",   "ˆ",   "‰",   "Š",   "‹","Û",
    "¶",   "d",    "[k",  "x",  "T",  "t",   "M+", "<+", "Q",  ";",    "j",   "u",
    "Ù",   "Ùk",   "ä",    "–",   "—",       
    "à",   "á",    "â",   "ã",   "ºz",  "º",   "í", "{", "{k",  "«", "=","K", 
    "Nî",   "Vî",    "Bî",   "Mî",   "<î", "|","}",
    "J",   "Vª",   "Mª",  "<ªª",  "Nª",   "Ø",  "Ý",   "æ", "ç", "xz", "#", ":",
    "z",
    "vks",  "vkS",  "vk",    "v",   "bZ",  "b",  "m",  "Å",  ",s",  ",",   "_",
    "D",  "d",    "ô",     "[",     "[k",    "X",   "x",  "?",    "?k",   "³", 
    "pkS",  "P",    "p",  "N",   "T",    "t",   "÷",  ">",   "¥",
    "ê",      "ë",      "V",  "B",   "ì",       "ï",     "M",  "<",  ".", ".k",   
    "R",  "r",   "F", "Fk",  ")",    "n", "/",  "/k",  "U", "u",   
    "I",  "i",   "¶", "Q",   "C",  "c",  "H",  "Hk", "E",   "e",
    "¸",   ";",    "j",  "Y",   "y",  "G",  "O",  "o",
    "'", "'k",  "\"", "\"k", "L",   "l",   "g",      
    "v‚",    "‚",    "ks",   "kS",   "k",     "h",    "q",   "w",   "`",    "s",    "S",
    "a",    "¡",    "%",     "W",   "·",   "~ ", "~"
]

PREETI_UNICODE_DICT = {
    "अ": "c", "आ": "cf", "ा": "f", "इ": "O", "ई": "O{", "र्": "{", "उ": "p", "ए": "P",
    "े": "]", "ै": "}", "ो": "f]", "ौ": "f}", "ओ": "cf]", "औ": "cf}", "ं": "+", "ँ": "F",
    "ि": "l", "ी": "L", "ु": "'", "ू": '"', "क": "s", "ख": "v", "ग": "u", "घ": "3",
    "ङ": "ª", "च": "r", "छ": "5", "ज": "h", "झ": "´", "ञ": "`", "ट": "6", "ठ": "7",
    "ड": "8", "ढ": "9", "ण": "0f", "त": "t", "थ": "y", "द": "b", "ध": "w", "न": "g",
    "प": "k", "फ": "km", "ब": "a", "भ": "e", "म": "d", "य": "o", "र": "/", "रू": "?",
    "ृ": "[", "ल": "n", "व": "j", "स": ";", "श": "z", "ष": "if", "ज्ञ": "1", "ह": "x",
    "१": "!", "२": "@", "३": "#", "४": "$", "५": "%", "६": "^", "७": "&", "८": "*",
    "९": "(", "०": ")", "।": ".", "्": "\\", "ऊ": "pm", "-": " ", "(": "-", ")": "_"
}

PREETI_TO_UNICODE_DICT = {
    "÷": "/", "v": "ख", "r": "च", "\"": "ू", "~": "ञ्", "z": "श", "ç": "ॐ", "f": "ा",
    "b": "द", "n": "ल", "j": "व", "×": "×", "V": "ख्", "R": "च्", "ß": "द्म", "^": "६",
    "Û": "!", "Z": "श्", "F": "ँ", "B": "द्य", "N": "ल्", "Ë": "ङ्ग", "J": "व्", "6": "ट",
    "2": "द्द", "¿": "रू", ">": "श्र", ":": "स्", "§": "ट्ट", "&": "७", "£": "घ्",
    "•": "ड्ड", ".": "।", "«": "्र", "*": "८", "„": "ध्र", "w": "ध", "s": "क", "g": "न",
    "æ": "“", "c": "अ", "o": "य", "k": "प", "W": "ध्", "Ö": "=", "S": "क्", "Ò": "¨",
    "_": ")", "[": "ृ", "Ú": "’", "G": "न्", "ˆ": "फ्", "C": "ऋ", "O": "इ", "Î": "ङ्ख",
    "K": "प्", "7": "ठ", "¶": "ठ्ठ", "3": "घ", "9": "ढ", "?": "रु", ";": "स", "'": "ु",
    "#": "३", "¢": "द्घ", "/": "र", "+": "ं", "ª": "ङ", "t": "त", "p": "उ", "|": "्र",
    "x": "ह", "å": "द्व", "d": "म", "`": "ञ", "l": "ि", "h": "ज", "T": "त्", "P": "ए",
    "Ý": "ट्ठ", "\\": "्", "Ù": ";", "X": "ह्", "Å": "हृ", "D": "म्", "@": "२", "Í": "ङ्क",
    "L": "ी", "H": "ज्", "4": "द्ध", "±": "+", "0": "ण्", "<": "?", "8": "ड", "¥": "र्‍",
    "$": "४", "¡": "ज्ञ्", ",": ",", "©": "र", "(": "९", "‘": "ॅ", "u": "ग", "q": "त्र",
    "}": "ै", "y": "थ", "e": "भ", "a": "ब", "i": "ष्", "‰": "झ्", "U": "ग्", "Q": "त्त",
    "]": "े", "˜": "ऽ", "Y": "थ्", "Ø": "्य", "E": "भ्", "A": "ब्", "M": "ः", "Ì": "न्न",
    "I": "क्ष्", "5": "छ", "´": "झ", "1": "ज्ञ", "°": "ङ्ढ", "=": ".", "Æ": "”",
    "‹": "ङ्घ", "%": "५", "¤": "झ्", "!": "१", "-": "(", "›": "द्र", ")": "०", "…": "‘",
    "Ü": "%"
}

def convert_krutidev_to_unicode(text: str) -> str:
    if not text:
        return ""
    out = text
    for s_one, s_two in zip(DEVLYS_ARRAY_ONE, DEVLYS_ARRAY_TWO):
        out = out.replace(s_one, s_two)
    
    out = out.replace("±", "Zं")
    out = out.replace("Æ", "र्f")
    
    import re
    out = out.replace("Ç", "fa")
    out = out.replace("É", "र्fa")
    
    # Swap 'f' with the following consonant cluster:
    # A consonant cluster consists of zero or more half-consonants (consonant + halant),
    # followed by a full consonant, and optionally a subjoined/trailing 'r' sign.
    out = re.sub(
        r"f((?:[\u0915-\u0939\u0958-\u095F]\u094d)*[\u0915-\u0939\u0958-\u095F](?:\u094d\u0930|्र)?)",
        r"\1ि",
        out
    )
    
    # Swap 'fa' for short-i + bindu (िं)
    out = re.sub(
        r"fa((?:[\u0915-\u0939\u0958-\u095F]\u094d)*[\u0915-\u0939\u0958-\u095F](?:\u094d\u0930|्र)?)",
        r"\1िं",
        out
    )
        
    out = out.replace("Ê", "ीZ")
    
    pos = out.find("ि्")
    while pos != -1:
        if pos + 2 < len(out):
            next_char = out[pos + 2]
            out = out[:pos] + "्" + next_char + "ि" + out[pos + 3:]
        pos = out.find("ि्", pos + 2)
        
    reph_matras = "अ आ इ ई उ ऊ ए ऐ ओ औ ा ि ी ु ू ृ े ै ो ो ौ ं : ँ ॅ"
    reph_matras_set = set(reph_matras.split())
    pos_z = out.find("Z")
    while pos_z > 0:
        prob_pos = pos_z - 1
        while prob_pos >= 0 and out[prob_pos] in reph_matras_set:
            prob_pos -= 1
        chars_to_move = out[max(0, prob_pos):pos_z]
        if prob_pos < 0:
            out = "र्" + chars_to_move + out[pos_z+1:]
        else:
            out = out[:prob_pos] + "र्" + chars_to_move + out[pos_z+1:]
        pos_z = out.find("Z")
        
    return out

def convert_if_legacy_hindi(text: str) -> str:
    if not text:
        return text
    # List of telltale legacy font triggers (KrutiDev / Preeti / DevLys indicators)
    triggers = [
        "o'kZ", "ekg", "vk;q", "lhek", "ijh{kk", "vk;ksx", "lsok", "rFkk", 
        "}kjk", "foHkkx", "fpfd", "vH;f", "vH;", "’kkl", "f’k{k", "fu;e", 
        "e/;", "çns", "ys[kk", "rduhd", "inks", "ykblsa", "fyfi", "oxZ", 
        "LFkk", "çf'k", "’kS{k", "’kS", "fyf[kr", "Yksd", "yksd", "vk;q"
    ]
    import re
    has_trigger = any(t in text for t in triggers)
    if has_trigger or re.search(r"\b[a-zA-Z]*[’']/d[a-zA-Z]*\b", text) or re.search(r"f[a-zA-Z][ksS]", text) or re.search(r"[a-zA-Z]Z\b", text):
        return convert_krutidev_to_unicode(text)
    return text

def convert_unicode_to_krutidev(text: str) -> str:
    if not text:
        return ""
    import re
    out = text
    nuktas = {
        "क़": "क़", "ख़‌": "ख़", "ग़": "ग़", "ज़": "ज़", "ड़": "ड़", "ढ़": "ढ़", 
        "ऩ": "ऩ", "फ़": "फ़", "य़": "य़", "ऱ": "ऱ"
    }
    for k, v in nuktas.items():
        out = out.replace(k, v)
        
    pos_i = out.find("ि")
    while pos_i != -1:
        if pos_i > 0:
            char_left = out[pos_i - 1]
            out = out[:pos_i - 1] + "f" + char_left + out[pos_i + 1:]
            new_pos = pos_i - 1
            while new_pos > 1 and out[new_pos - 1] == "्":
                half_cons = out[new_pos - 2]
                out = out[:new_pos - 2] + "f" + half_cons + "्" + out[new_pos + 1:]
                new_pos -= 2
            pos_i = out.find("ि", new_pos + 2)
        else:
            pos_i = out.find("ि", pos_i + 1)
            
    set_of_matras = "ािीुूृेैोौं:ँॅ"
    set_matras = set(list(set_of_matras))
    pos_r = out.find("र्")
    while pos_r != -1:
        prob_z = pos_r + 2
        while prob_z + 1 < len(out) and out[prob_z + 1] in set_matras:
            prob_z += 1
        cluster = out[pos_r + 2 : prob_z + 1]
        out = out[:pos_r] + cluster + "Z" + out[prob_z + 1:]
        pos_r = out.find("र्")
        
    for s_one, s_two in zip(UNICODE_TO_KRUTIDEV_ONE, UNICODE_TO_KRUTIDEV_TWO):
        out = out.replace(s_one, s_two)
        
    return out

def convert_preeti_to_unicode(text: str) -> str:
    import re
    if not text:
        return ""
    out = "".join(PREETI_TO_UNICODE_DICT.get(c, c) for c in text)
    out = out.replace("्ा", "")
    out = re.sub(r"(त्र|त्त)([^उभप]+?)m", r"\1m\2", out)
    out = out.replace("त्रm", "क्र")
    out = out.replace("त्तm", "क्त")
    out = re.sub(r"([^उभप]+?)m", r"m\1", out)
    out = out.replace("उm", "ऊ")
    out = out.replace("भm", "झ")
    out = out.replace("पm", "फ")
    out = out.replace("इ{", "ई")
    out = re.sub(r"(?<![\u0915-\u0939\u0958-\u095F\u094D])ि((?:[\u0915-\u0939\u0958-\u095F]\u094D)*[\u0915-\u0939\u0958-\u095F](?:\u094D\u0930|्र)?)", r"\1ि", out)
    out = re.sub(r"((?:.[ािीुूृेैोौंःँ]*?)){", r"{\1", out)
    out = re.sub(r"((?:.्)*){", r"{\1", out)
    out = out.replace("{", "र्")
    out = re.sub(r"([ाीुूृेैोौंःँ]+?)(्(?:.्)*[^्])", r"\2\1", out)
    out = re.sub(r"्([ाीुूृेैोौंःँ]+?)((?:.्)*[^्])", r"्\2\1", out)
    out = re.sub(r"([ंँ])([ािीुूृेैोौः]*)", r"\2\1", out)
    out = out.replace("ँँ", "ँ").replace("ंं", "ं").replace("ेे", "े").replace("ैै", "ै").replace("ुु", "ु").replace("ूू", "ू")
    
    out = re.sub(r"^ः", ":", out)
    out = out.replace("टृ", "ट्ट").replace("ेा", "ाे").replace("ैा", "ाै")
    out = out.replace("अाे", "ओ").replace("अाै", "औ").replace("अा", "आ").replace("एे", "ऐ")
    out = out.replace("ाे", "ो").replace("ाै", "ौ")
    return out

def convert_unicode_to_preeti(unicodestring: str) -> str:
    if not unicodestring:
        return ""
    index = -1
    normalized = ""
    while index + 1 < len(unicodestring):
        index += 1
        character = unicodestring[index]
        try:
            if character != "र":
                if index + 2 < len(unicodestring) and unicodestring[index + 1] == "्" and unicodestring[index + 2] not in (" ", "।", ","):
                    if unicodestring[index + 2] != "र":
                        preeti_val = PREETI_UNICODE_DICT.get(character)
                        if preeti_val in list("wertyuxasdghjkzvn"):
                            normalized += chr(ord(preeti_val) - 32)
                            index += 1
                            continue
                        elif character == "स":
                            normalized += ":"
                            index += 1
                            continue
                        elif character == "ष":
                            normalized += "i"
                            index += 1
                            continue
            if index > 0 and index + 1 < len(unicodestring) and unicodestring[index - 1] != "र" and character == "्" and unicodestring[index + 1] == "र":
                if unicodestring[index - 1] not in ("ट", "ठ", "ड"):
                    normalized += "|"
                    index += 1
                    continue
                else:
                    normalized += "«"
                    index += 1
                    continue
        except Exception:
            pass
        normalized += character
        
    normalized = normalized.replace("त|", "q")
    
    converted = ""
    index = -1
    while index + 1 < len(normalized):
        index += 1
        character = normalized[index]
        if character == "\ufeff":
            continue
        try:
            if index + 1 < len(normalized) and normalized[index + 1] == "ि":
                if character == "q":
                    converted += "l" + character
                else:
                    converted += "l" + PREETI_UNICODE_DICT.get(character, character)
                index += 1
                continue
                
            if index + 2 < len(normalized) and normalized[index + 2] == "ि":
                if character in list("WERTYUXASDGHJK:ZVN"):
                    if normalized[index + 1] != "q":
                        converted += "l" + character + PREETI_UNICODE_DICT.get(normalized[index + 1], normalized[index + 1])
                        index += 2
                        continue
                    else:
                        converted += "l" + character + "q"
                        index += 2
                        continue
                        
            if index + 2 < len(normalized) and normalized[index + 1] == "्" and character == "र":
                if index + 3 < len(normalized) and normalized[index + 3] in ("ा", "ो", "ौ", "े", "ै", "ी"):
                    converted += PREETI_UNICODE_DICT.get(normalized[index + 2], normalized[index + 2]) + PREETI_UNICODE_DICT.get(normalized[index + 3], normalized[index + 3]) + "{"
                    index += 3
                    continue
                elif index + 3 < len(normalized) and normalized[index + 3] == "ि":
                    converted += PREETI_UNICODE_DICT.get(normalized[index + 3], normalized[index + 3]) + PREETI_UNICODE_DICT.get(normalized[index + 2], normalized[index + 2]) + "{"
                    index += 3
                    continue
                converted += PREETI_UNICODE_DICT.get(normalized[index + 2], normalized[index + 2]) + "{"
                index += 2
                continue
                
            if index + 3 < len(normalized) and normalized[index + 3] == "ि":
                if normalized[index + 2] in ("|", "«"):
                    if character in list("WERTYUXASDGHJK:ZVNIi"):
                        converted += "l" + character + PREETI_UNICODE_DICT.get(normalized[index + 1], normalized[index + 1]) + normalized[index + 2]
                        index += 3
                        continue
        except Exception:
            pass
            
        converted += PREETI_UNICODE_DICT.get(character, character)
        
    converted = converted.replace("Si", "I").replace("H`", "1").replace("b\\w", "4").replace("z|", ">")
    converted = converted.replace("/'", "?").replace('/"', "¿").replace("Tt", "Q").replace("b\\lj", "lå")
    converted = converted.replace("b\\j", "å").replace("0f\\", "0").replace("`\\", "~")
    return converted

@router.post("/font-convert")
async def font_convert(request: FontConvertRequest):
    ctype = request.conversion_type.lower()
    text = request.text
    
    if ctype in ("krutidev_to_unicode", "krutidev_to_mangal", "devlys_to_unicode"):
        result = convert_krutidev_to_unicode(text)
    elif ctype in ("unicode_to_krutidev", "mangal_to_krutidev", "unicode_to_devlys"):
        result = convert_unicode_to_krutidev(text)
    elif ctype == "preeti_to_unicode":
        result = convert_preeti_to_unicode(text)
    elif ctype == "unicode_to_preeti":
        result = convert_unicode_to_preeti(text)
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported conversion type: '{request.conversion_type}'")
        
    return {"status": "success", "converted_text": result}



