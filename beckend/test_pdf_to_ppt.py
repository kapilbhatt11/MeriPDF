import asyncio
import os
import sys
import io
import fitz
from pptx import Presentation
from pptx.util import Pt

# Directly test PyMuPDF page rendering to PPTX slide (Visual Replica mode logic)
def test_standalone_conversion():
    sample_pdf_path = os.path.join(os.path.dirname(__file__), "sample.pdf")
    if not os.path.exists(sample_pdf_path):
        print("❌ sample.pdf not found!")
        return

    print("Opening sample PDF...")
    doc = fitz.open(sample_pdf_path)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    first_page = doc.load_page(0)
    prs.slide_width = int(Pt(first_page.rect.width))
    prs.slide_height = int(Pt(first_page.rect.height))

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=300)
        img_path = os.path.join(os.path.dirname(__file__), f"test_page_{page_num}.png")
        pix.save(img_path)

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

        if os.path.exists(img_path):
            os.remove(img_path)

    out_pptx = os.path.join(os.path.dirname(__file__), "test_standalone_out.pptx")
    prs.save(out_pptx)
    print(f"✅ Standalone PDF-to-PPT Conversion Successful! Created {out_pptx} with {len(prs.slides)} slides.")

if __name__ == "__main__":
    test_standalone_conversion()
